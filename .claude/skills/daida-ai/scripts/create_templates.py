#!/usr/bin/env python3
"""PPTXテンプレート生成スクリプト

tech / casual / formal の3種類のテンプレートを生成する。
各テンプレートには6つのスライドレイアウトを含み、
slide_builder.py が期待するプレースホルダー構成に一致させる。

NOTE: python-pptxではスライドマスターの完全なカスタマイズが困難なため、
このスクリプトはデフォルトテンプレートをベースにスタイル（色・フォント）を
適用したサンプルスライドを含むテンプレートを生成する。
本番クオリティにはLibreOffice Impressでの追加調整を推奨。
"""

import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

SCRIPT_DIR = Path(__file__).resolve().parent
TEMPLATE_DIR = SCRIPT_DIR.parent / "assets" / "templates"

# テンプレート定義: 配色とフォント設定
TEMPLATES = {
    "tech": {
        "bg_dark": RGBColor(0x1E, 0x1E, 0x2E),      # ダーク背景
        "bg_section": RGBColor(0x28, 0x28, 0x3C),    # セクション背景
        "title_color": RGBColor(0x89, 0xB4, 0xFA),   # 明るい青
        "text_color": RGBColor(0xCD, 0xD6, 0xF4),    # ライトグレー
        "accent_color": RGBColor(0xA6, 0xE3, 0xA1),  # グリーン
        "subtitle_color": RGBColor(0x94, 0x99, 0xB7), # サブグレー
    },
    "casual": {
        "bg_dark": RGBColor(0xFF, 0xFF, 0xFF),        # ホワイト背景
        "bg_section": RGBColor(0xF0, 0xF4, 0xFF),    # 薄い青背景
        "title_color": RGBColor(0x2D, 0x3A, 0x8C),   # ネイビー
        "text_color": RGBColor(0x33, 0x33, 0x33),    # ダークグレー
        "accent_color": RGBColor(0xFF, 0x6B, 0x35),  # オレンジ
        "subtitle_color": RGBColor(0x66, 0x66, 0x66), # ミッドグレー
    },
    "formal": {
        "bg_dark": RGBColor(0xF8, 0xF8, 0xF8),       # オフホワイト
        "bg_section": RGBColor(0x00, 0x2B, 0x5C),    # ネイビー背景
        "title_color": RGBColor(0x00, 0x2B, 0x5C),   # ネイビー
        "text_color": RGBColor(0x33, 0x33, 0x33),    # ダークグレー
        "accent_color": RGBColor(0xC4, 0x96, 0x3C),  # ゴールド
        "subtitle_color": RGBColor(0x55, 0x55, 0x55), # グレー
    },
}


def _set_slide_bg(slide, color: RGBColor):
    """スライドの背景色を設定する"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_styled_text(text_frame, text, color, size_pt, bold=False, alignment=None):
    """スタイル付きテキストを追加する"""
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = text
    p.font.color.rgb = color
    p.font.size = Pt(size_pt)
    p.font.bold = bold
    if alignment:
        p.alignment = alignment


def create_template(name: str, colors: dict) -> Path:
    """指定されたスタイルでテンプレートPPTXを生成する"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Layout 0: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_bg(slide, colors["bg_dark"])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in phs:
        _add_styled_text(
            phs[0].text_frame, "プレゼンテーションタイトル",
            colors["title_color"], 44, bold=True, alignment=PP_ALIGN.CENTER
        )
    if 1 in phs:
        _add_styled_text(
            phs[1].text_frame, "サブタイトル / 登壇者名",
            colors["subtitle_color"], 24, alignment=PP_ALIGN.CENTER
        )

    # Layout 1: Title and Content
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    _set_slide_bg(slide, colors["bg_dark"])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in phs:
        _add_styled_text(
            phs[0].text_frame, "スライドタイトル",
            colors["title_color"], 36, bold=True
        )
    if 1 in phs:
        tf = phs[1].text_frame
        tf.clear()
        for i, item in enumerate(["ポイント1: 具体的な内容", "ポイント2: データや数値", "ポイント3: まとめ"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = item
            p.font.color.rgb = colors["text_color"]
            p.font.size = Pt(24)

    # Layout 2: Section Header
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    _set_slide_bg(slide, colors["bg_section"])
    if slide.shapes.title:
        _add_styled_text(
            slide.shapes.title.text_frame, "セクションタイトル",
            colors["accent_color"] if name == "tech" else colors["title_color"],
            40, bold=True, alignment=PP_ALIGN.CENTER
        )

    # Layout 3: Two Content
    slide = prs.slides.add_slide(prs.slide_layouts[3])
    _set_slide_bg(slide, colors["bg_dark"])
    phs = {ph.placeholder_format.idx: ph for ph in slide.placeholders}
    if 0 in phs:
        _add_styled_text(
            phs[0].text_frame, "比較タイトル",
            colors["title_color"], 36, bold=True
        )
    for idx in [1, 2]:
        if idx in phs:
            tf = phs[idx].text_frame
            tf.clear()
            side = "左" if idx == 1 else "右"
            p = tf.paragraphs[0]
            p.text = f"{side}カラム見出し"
            p.font.color.rgb = colors["accent_color"]
            p.font.size = Pt(22)
            p.font.bold = True
            for item in ["項目A", "項目B"]:
                p2 = tf.add_paragraph()
                p2.text = item
                p2.font.color.rgb = colors["text_color"]
                p2.font.size = Pt(20)

    # Layout 5: Title Only
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _set_slide_bg(slide, colors["bg_dark"])
    if slide.shapes.title:
        _add_styled_text(
            slide.shapes.title.text_frame, "図・画像用スライド",
            colors["title_color"], 36, bold=True
        )

    # Layout 6: Blank
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, colors["bg_dark"])

    # テンプレートとして使うため、全スライドを削除（レイアウト定義のみ残す）
    # python-pptxにはスライド削除APIがないため、XML操作とパーツ削除で対応
    for slide in list(prs.slides):
        rId = slide.part.partname
        # プレゼンテーションパーツからスライドパーツへの関連を削除
        for rel in prs.part.rels.values():
            if rel.target_part is slide.part:
                prs.part.drop_rel(rel.rId)
                break
    # sldIdLst からも参照を除去
    slide_list = prs.slides._sldIdLst
    for sldId in list(slide_list):
        slide_list.remove(sldId)

    # Save
    TEMPLATE_DIR.mkdir(parents=True, exist_ok=True)
    output_path = TEMPLATE_DIR / f"{name}.pptx"
    prs.save(str(output_path))
    return output_path


def main():
    print("=== Creating PPTX templates ===")
    for name, colors in TEMPLATES.items():
        path = create_template(name, colors)
        print(f"  {name}: {path}")
    print("Done.")


if __name__ == "__main__":
    main()
