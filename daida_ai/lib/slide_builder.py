"""スライド仕様(SlideSpec) → python-pptx Presentationへの変換"""

from __future__ import annotations

from pathlib import Path

from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt, Emu

from daida_ai.lib.slide_spec import SlideSpec, Slide, TwoColumnContent

# 画像配置のマージン定数 (EMU)
_IMG_MARGIN = Emu(457200)         # 左右・下マージン
_IMG_TOP = Emu(1600200)           # タイトル下の画像開始位置
_IMG_TOP_BLANK = Emu(457200)     # blankレイアウト時の画像開始位置

# デフォルトテンプレートのレイアウトインデックス
_DEFAULT_LAYOUT_IDX = {
    "title_slide": 0,
    "title_and_content": 1,
    "section_header": 2,
    "two_content": 3,
    "title_only": 5,
    "blank": 6,
}


def _find_layout(prs, layout_name: str):
    """レイアウトをインデックスまたは名前で取得する。"""
    idx = _DEFAULT_LAYOUT_IDX.get(layout_name)
    if idx is not None:
        try:
            return prs.slide_layouts[idx]
        except IndexError:
            pass

    name_map = {
        "title_slide": "Title Slide",
        "title_and_content": "Title and Content",
        "section_header": "Section Header",
        "two_content": "Two Content",
        "title_only": "Title Only",
        "blank": "Blank",
    }
    target_name = name_map.get(layout_name, layout_name)
    for layout in prs.slide_layouts:
        if layout.name == target_name:
            return layout

    return prs.slide_layouts[0]


def _get_placeholders(slide) -> dict:
    """スライドのplaceholdersをidx→placeholder の辞書として返す。"""
    return {ph.placeholder_format.idx: ph for ph in slide.placeholders}


def _set_notes(slide, text: str | None) -> None:
    """スライドにスピーカーノートを設定する。"""
    if text:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text


def _get_slide_size(slide):
    """スライドの幅・高さを Presentation から取得する。"""
    prs = slide.part.package.presentation_part.presentation
    return int(prs.slide_width), int(prs.slide_height)


def _calc_image_area(slide_w: int, slide_h: int, *, is_blank: bool = False):
    """画像配置可能領域の最大幅・最大高さ・開始位置を計算する。"""
    margin = int(_IMG_MARGIN)
    top = int(_IMG_TOP_BLANK) if is_blank else int(_IMG_TOP)
    max_w = slide_w - 2 * margin
    max_h = slide_h - top - margin
    return max_w, max_h, top


def _convert_svg(svg_path: Path, original_path: str, *, is_blank: bool = False) -> Path:
    """SVGをPNGに変換し、一時ファイルのPathを返す。

    変換前にフォントサイズのバリデーションを実行し、
    PPTX上で12pt未満になるテキストがあれば警告を出す。

    Raises:
        FileNotFoundError: cairosvg未インストール or 変換失敗時
    """
    import warnings

    from daida_ai.lib.svg_convert import (
        convert_svg_to_png,
        SVGConversionError,
        validate_svg_font_sizes,
    )

    # フォントサイズ検証
    try:
        svg_content = svg_path.read_text(encoding="utf-8")
        violations = validate_svg_font_sizes(svg_content, is_blank=is_blank)
        for v in violations:
            warnings.warn(
                f"SVG font too small in {original_path}: {v}",
                UserWarning,
                stacklevel=2,
            )
    except Exception:
        pass  # バリデーション失敗は変換をブロックしない

    try:
        png_path = convert_svg_to_png(str(svg_path))
        return Path(png_path)
    except SVGConversionError as e:
        raise FileNotFoundError(
            f"SVG conversion failed: {original_path} ({e})"
        ) from e


def _insert_image(slide, image_path: str, *, is_blank: bool = False, base_dir: Path | None = None) -> None:
    """スライドに画像を挿入する。アスペクト比を維持してコンテンツ領域に収める。

    スライドの実際のサイズから配置領域を動的に計算するため、
    4:3、16:9、カスタムサイズのテンプレートいずれにも対応する。

    Raises:
        FileNotFoundError: 画像ファイルが存在しないか、読み取れない
        ValueError: 画像パスがbase_dirの外を参照している（パストラバーサル防止）
    """
    path = Path(image_path)
    if not path.is_absolute() and base_dir is not None:
        path = (base_dir / path).resolve()
        # パストラバーサル防止: 解決後のパスがbase_dir配下にあることを検証
        base_resolved = base_dir.resolve()
        if not path.is_relative_to(base_resolved):
            raise ValueError(
                f"Image path escapes base directory: {image_path}"
            )
    if not path.exists():
        raise FileNotFoundError(f"Image not found: {image_path}")

    # SVG自動変換: .svg ファイルを検出し、PNGに変換してから挿入
    temp_png = None
    if path.suffix.lower() == ".svg":
        path = _convert_svg(path, image_path, is_blank=is_blank)
        temp_png = path

    try:
        try:
            with PILImage.open(str(path)) as img:
                img_w, img_h = img.size
        except (OSError, SyntaxError) as e:
            raise FileNotFoundError(f"Invalid image file: {image_path} ({e})") from e

        slide_w, slide_h = _get_slide_size(slide)
        max_w, max_h, top = _calc_image_area(slide_w, slide_h, is_blank=is_blank)

        # アスペクト比を維持してフィットさせる
        scale = min(max_w / img_w, max_h / img_h)
        width = Emu(int(img_w * scale))
        height = Emu(int(img_h * scale))

        # 水平中央, 垂直はコンテンツ領域内で中央
        left = Emu((slide_w - int(width)) // 2)
        top_pos = Emu(top + (max_h - int(height)) // 2)

        return slide.shapes.add_picture(str(path), left, top_pos, width, height)
    finally:
        if temp_png is not None:
            temp_png.unlink(missing_ok=True)


def _send_to_back(slide, shape) -> None:
    """シェイプをspTree内の最前面から最背面に移動する。"""
    spTree = slide.element.find(
        ".//{http://schemas.openxmlformats.org/presentationml/2006/main}cSld/"
        "{http://schemas.openxmlformats.org/presentationml/2006/main}spTree"
    )
    if spTree is None:
        return
    sp_el = shape._element
    spTree.remove(sp_el)
    # spTree の最初の子要素（nvGrpSpPr, grpSpPr）の後に挿入
    insert_idx = 0
    for i, child in enumerate(spTree):
        tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
        if tag in ("nvGrpSpPr", "grpSpPr"):
            insert_idx = i + 1
        else:
            break
    spTree.insert(insert_idx, sp_el)


def _add_title_slide(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "title_slide")
    slide = prs.slides.add_slide(layout)
    phs = _get_placeholders(slide)
    if 0 in phs:
        phs[0].text = spec.title
    if spec.subtitle and 1 in phs:
        phs[1].text = spec.subtitle
    if spec.image:
        pic = _insert_image(slide, spec.image, base_dir=base_dir)
        _send_to_back(slide, pic)
    _set_notes(slide, spec.note)


def _add_section_header(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "section_header")
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.title
    if spec.image:
        pic = _insert_image(slide, spec.image, base_dir=base_dir)
        _send_to_back(slide, pic)
    _set_notes(slide, spec.note)


def _add_title_and_content(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "title_and_content")
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.title

    phs = _get_placeholders(slide)
    if spec.body and 1 in phs:
        tf = phs[1].text_frame
        tf.clear()
        for i, item in enumerate(spec.body):
            if i == 0:
                tf.paragraphs[0].text = item
            else:
                p = tf.add_paragraph()
                p.text = item

    if spec.image:
        pic = _insert_image(slide, spec.image, base_dir=base_dir)
        _send_to_back(slide, pic)
    _set_notes(slide, spec.note)


def _add_two_content(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "two_content")
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.title

    phs = _get_placeholders(slide)

    if spec.left and 1 in phs:
        tf = phs[1].text_frame
        tf.clear()
        tf.paragraphs[0].text = spec.left.heading
        tf.paragraphs[0].font.bold = True
        for item in spec.left.body:
            p = tf.add_paragraph()
            p.text = item

    if spec.right and 2 in phs:
        tf = phs[2].text_frame
        tf.clear()
        tf.paragraphs[0].text = spec.right.heading
        tf.paragraphs[0].font.bold = True
        for item in spec.right.body:
            p = tf.add_paragraph()
            p.text = item

    if spec.image:
        pic = _insert_image(slide, spec.image, base_dir=base_dir)
        _send_to_back(slide, pic)
    _set_notes(slide, spec.note)


def _add_title_only(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "title_only")
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = spec.title
    if spec.image:
        pic = _insert_image(slide, spec.image, base_dir=base_dir)
        _send_to_back(slide, pic)
    _set_notes(slide, spec.note)


def _add_blank(prs, spec: Slide, base_dir: Path | None = None) -> None:
    layout = _find_layout(prs, "blank")
    slide = prs.slides.add_slide(layout)
    if spec.image:
        _insert_image(slide, spec.image, is_blank=True, base_dir=base_dir)
    _set_notes(slide, spec.note)


_BUILDERS = {
    "title_slide": _add_title_slide,
    "section_header": _add_section_header,
    "title_and_content": _add_title_and_content,
    "two_content": _add_two_content,
    "title_only": _add_title_only,
    "blank": _add_blank,
}


def build_presentation(
    spec: SlideSpec,
    template_path: str | None = None,
    base_dir: Path | None = None,
):
    """SlideSpecからpython-pptx Presentationを生成する。

    Args:
        spec: スライド仕様
        template_path: カスタムテンプレートPPTXのパス（Noneでデフォルト）
        base_dir: 画像の相対パスを解決するベースディレクトリ

    Returns:
        python-pptx Presentation オブジェクト
    """
    prs = Presentation(template_path)

    for slide_spec in spec.slides:
        builder = _BUILDERS.get(slide_spec.layout)
        if builder:
            builder(prs, slide_spec, base_dir=base_dir)

    return prs
