"""スライドショー自動再生設定

PPTXにauto-advance transition と 音声auto-play animation を設定し、
スライドショー開始から終了まで完全自動で走るようにする。
"""

from __future__ import annotations

from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.opc.package import RT

_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_nsmap = {"p": _P_NS, "a": _A_NS, "r": _R_NS}

# MP3フレームベースの簡易デュレーション推定用
# MPEG-1 Layer III ビットレートテーブル (kbps)
_MPEG1_L3_BITRATE = {
    0b0001: 32, 0b0010: 40, 0b0011: 48, 0b0100: 56,
    0b0101: 64, 0b0110: 80, 0b0111: 96, 0b1000: 112,
    0b1001: 128, 0b1010: 160, 0b1011: 192, 0b1100: 224,
    0b1101: 256, 0b1110: 320,
}
# MPEG-2/2.5 Layer III ビットレートテーブル (kbps)
_MPEG2_L3_BITRATE = {
    0b0001: 8, 0b0010: 16, 0b0011: 24, 0b0100: 32,
    0b0101: 40, 0b0110: 48, 0b0111: 56, 0b1000: 64,
    0b1001: 80, 0b1010: 96, 0b1011: 112, 0b1100: 128,
    0b1101: 144, 0b1110: 160,
}


_DEFAULT_UNMEASURABLE_DURATION_MS = 30000

# ノートベースタイミング推定用
# 日本語の平均的な発話速度: 約300文字/分 = 5文字/秒
_CHARS_PER_SECOND = 5.0
_MIN_NOTE_DURATION_MS = 3000
_MAX_NOTE_DURATION_MS = 120000

# ECMA-376 CT_Slide 子要素の正規順序
# cSld, clrMapOvr?, transition?, timing?, hf?, extLst?
_CT_SLIDE_ORDER = ["cSld", "clrMapOvr", "transition", "timing", "hf", "extLst"]


def configure_slideshow(
    input_path: Path,
    output_path: Path,
    *,
    silent_duration_ms: int = 3000,
    audio_buffer_ms: int = 1000,
    unmeasurable_duration_ms: int = _DEFAULT_UNMEASURABLE_DURATION_MS,
) -> None:
    """PPTXにスライドショー自動再生設定を追加する。

    Args:
        input_path: 入力PPTXファイルパス
        output_path: 出力PPTXファイルパス
        silent_duration_ms: 音声なしスライドの表示時間（ミリ秒）
        audio_buffer_ms: 音声付きスライドの再生完了後の余白（ミリ秒）
        unmeasurable_duration_ms: デュレーション計測不能時のフォールバック（ミリ秒）
    """
    prs = Presentation(str(input_path))

    for slide in prs.slides:
        has_audio_shapes = bool(_find_audio_shape_ids(slide))
        audio_duration = _get_audio_duration_ms(slide)

        if has_audio_shapes:
            raw_dur = audio_duration if audio_duration > 0 else unmeasurable_duration_ms
            main_seq_dur_ms = max(raw_dur, 1)

            # 既存アニメーション終了時刻を考慮して advance_ms を調整する。
            # advTm は mainSeq.dur (= max(音声長, 既存アニメーション長)) + buffer 以上にする。
            # これにより PowerPoint でも LibreOffice でも既存アニメーションが
            # 完了する前にスライドが進むことを防ぐ。
            existing_end = _get_existing_anim_end_ms(slide)
            effective_dur = max(main_seq_dur_ms, existing_end)
            advance_ms = effective_dur + audio_buffer_ms
        else:
            note_duration = _estimate_note_duration_ms(
                slide, min_ms=silent_duration_ms
            )
            if note_duration > 0:
                advance_ms = note_duration
            else:
                advance_ms = silent_duration_ms
            # 音声なしスライドでも既存アニメーション長を考慮する
            existing_end = _get_existing_anim_end_ms(slide)
            if existing_end > 0:
                advance_ms = max(advance_ms, existing_end + audio_buffer_ms)

        # ECMA-376 CT_Slide: transition? は timing? より前に来る必要がある
        # transitionを先に設定してからtimingを追加する
        _set_auto_advance(slide, advance_ms)
        if has_audio_shapes:
            # LibreOffice は mainSeq dur="indefinite" のままだと
            # notifySlideAnimationsEnded() が呼ばれず advTm が発火しない。
            # 音声長を dur に設定することで音声終了時にアニメーションが終わり、
            # LibreOffice のページ自動送りが正しく動作する。
            _add_auto_play_animation(slide, main_seq_dur_ms=main_seq_dur_ms)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def _insert_slide_child(slide_elem: etree._Element, new_elem: etree._Element) -> None:
    """CT_Slideのスキーマ順序を守って子要素を挿入する。

    ECMA-376 CT_Slide: cSld, clrMapOvr?, transition?, timing?, hf?, extLst?
    new_elemのタグ名から正しい挿入位置を算出し、後続要素の直前にinsertする。
    """
    new_local = etree.QName(new_elem).localname
    if new_local not in _CT_SLIDE_ORDER:
        slide_elem.append(new_elem)
        return

    new_order = _CT_SLIDE_ORDER.index(new_local)

    # 自分より後ろの要素を探し、その直前に挿入する
    for child in slide_elem:
        child_local = etree.QName(child).localname
        if child_local in _CT_SLIDE_ORDER:
            child_order = _CT_SLIDE_ORDER.index(child_local)
            if child_order > new_order:
                child.addprevious(new_elem)
                return

    # 後続要素がなければ末尾に追加
    slide_elem.append(new_elem)


def _get_audio_duration_ms(slide) -> int:
    """スライドに埋め込まれた音声のデュレーションを推定する（ミリ秒）。

    音声がない場合は0を返す。
    複数の音声トラックがある場合は最長のデュレーションを返す
    （audioノードはdelay="0"で同時再生されるため、maxが正しい）。
    RT.MEDIAにはビデオも含まれるため、content_typeでaudioのみフィルタする。
    """
    media_rels = [r for r in slide.part.rels.values() if r.reltype == RT.MEDIA]
    if not media_rels:
        return 0

    max_duration = 0
    for rel in media_rels:
        if rel.is_external:
            continue
        try:
            part = rel.target_part
        except Exception:
            continue
        # ビデオ等の非音声メディアを除外
        if not part.content_type.startswith("audio/"):
            continue
        duration = _estimate_mp3_duration_ms(part.blob)
        if duration > max_duration:
            max_duration = duration
    return max_duration


def _estimate_mp3_duration_ms(data: bytes) -> int:
    """MP3バイナリからデュレーションをミリ秒で推定する。

    pydub依存を避け、ファイルサイズとビットレートから概算する。
    MP3フレームヘッダが見つからない場合（AAC/WAV等）は0を返す。
    """
    if len(data) < 10:
        return 0

    offset = 0
    # ID3タグをスキップ
    if data[:3] == b"ID3":
        if len(data) < 10:
            return 0
        size = (
            (data[6] & 0x7F) << 21
            | (data[7] & 0x7F) << 14
            | (data[8] & 0x7F) << 7
            | (data[9] & 0x7F)
        )
        offset = 10 + size

    # 最初のフレームヘッダを探す
    while offset < len(data) - 4:
        if data[offset] == 0xFF and (data[offset + 1] & 0xE0) == 0xE0:
            # MPEGバージョン判定: ビット4-3 of byte1
            # 11=MPEG1, 10=MPEG2, 00=MPEG2.5
            mpeg_version_bits = (data[offset + 1] >> 3) & 0x03
            is_mpeg1 = mpeg_version_bits == 0b11
            bitrate_table = _MPEG1_L3_BITRATE if is_mpeg1 else _MPEG2_L3_BITRATE
            default_bitrate = 128 if is_mpeg1 else 64

            bitrate_idx = (data[offset + 2] >> 4) & 0x0F
            bitrate_kbps = bitrate_table.get(bitrate_idx, default_bitrate)
            audio_bytes = len(data) - offset
            duration_s = (audio_bytes * 8) / (bitrate_kbps * 1000)
            return max(int(duration_s * 1000), 1)
        offset += 1

    # MP3フレームヘッダが見つからない（AAC/WAV等）→ 計測不能
    return 0


def _set_auto_advance(slide, advance_ms: int) -> None:
    """スライドに自動ページ送りを設定する。

    既存のtransition要素があれば属性を更新、なければ追加する。
    advClick="0" でクリックを無効化し、advTm でミリ秒指定。
    """
    slide_elem = slide.element
    trans = slide_elem.find(f"{{{_P_NS}}}transition")

    if trans is None:
        trans = etree.Element(f"{{{_P_NS}}}transition")
        _insert_slide_child(slide_elem, trans)

    trans.set("advClick", "0")
    trans.set("advTm", str(advance_ms))


def _add_auto_play_animation(slide, *, main_seq_dur_ms: int = 0) -> None:
    """スライドの音声シェイプに自動再生アニメーションを設定する。

    既存のアニメーション（テキスト・図形等）を保持しつつ、
    音声auto-playノードをmainSeqにマージする。

    Args:
        slide: python-pptxのスライドオブジェクト
        main_seq_dur_ms: mainSeqのdur属性をミリ秒で指定（0の場合は"indefinite"）。
            LibreOffice互換性のため音声長を渡すと、アニメーション終了が
            notifySlideAnimationsEnded() を確実に発火させる。

    PowerPointのアニメーション構造:
    p:timing > p:tnLst > p:par (tmRoot) > p:childTnLst > p:seq (mainSeq) >
    p:cTn > p:childTnLst > p:par > p:cTn > p:childTnLst > p:par > p:cTn >
    p:childTnLst > p:audio > p:cMediaNode > p:cTn + p:tgtEl
    """
    audio_shape_ids = _find_audio_shape_ids(slide)
    if not audio_shape_ids:
        return

    slide_elem = slide.element
    existing_timing = slide_elem.find(f"{{{_P_NS}}}timing")

    if existing_timing is not None:
        # 既存のtiming構造にaudioノードをマージ
        _merge_audio_into_timing(existing_timing, audio_shape_ids,
                                  main_seq_dur_ms=main_seq_dur_ms)
    else:
        # timing要素がない場合は新規作成
        timing_elem = _build_timing_xml(audio_shape_ids,
                                         main_seq_dur_ms=main_seq_dur_ms)
        _insert_slide_child(slide_elem, timing_elem)


def _merge_audio_into_timing(
    timing_elem: etree._Element,
    audio_shape_ids: list[int],
    *,
    main_seq_dur_ms: int = 0,
) -> None:
    """既存のp:timing構造に音声auto-playノードを追加する。

    mainSeqのchildTnLstに音声用p:parを追加する。
    既存のアニメーションはそのまま保持される。

    Args:
        timing_elem: 既存の p:timing 要素
        audio_shape_ids: 音声シェイプのIDリスト
        main_seq_dur_ms: 0の場合は "indefinite" のまま変更しない。
            正の値の場合は mainSeq の dur をその値（ms）に更新する。
    """
    # mainSeq (nodeType="mainSeq") を探す
    main_seq_ctn = timing_elem.find(
        f".//{{{_P_NS}}}cTn[@nodeType='mainSeq']"
    )
    if main_seq_ctn is None:
        return

    # LibreOffice互換: mainSeq の dur を更新する。
    # 比較対象は「以前のmainSeq.dur値（ステール可能性あり）」ではなく、
    # 「childTnLst内の明示的なdur属性を持つ非音声アニメーション」にする。
    # 音声ノードはdur属性を持たないため自然にフィルタされる。
    # これにより: (1)再実行時に短縮された音声長が反映される (2)既存非音声アニメーションが保護される
    if main_seq_dur_ms > 0:
        max_child = _get_max_child_animation_dur_ms(main_seq_ctn)
        main_seq_ctn.set("dur", str(max(main_seq_dur_ms, max_child)))

    child_tn_lst = main_seq_ctn.find(f"{{{_P_NS}}}childTnLst")
    if child_tn_lst is None:
        child_tn_lst = etree.SubElement(main_seq_ctn, f"{{{_P_NS}}}childTnLst")

    # 既存のメディアコマンドノードのspidを取得して重複を避ける
    existing_spids: set[int] = set()
    for spgt in timing_elem.iter(f"{{{_P_NS}}}spTgt"):
        parent = spgt.getparent()
        while parent is not None:
            if parent.tag in (f"{{{_P_NS}}}cmd", f"{{{_P_NS}}}audio"):
                spid = spgt.get("spid")
                if spid is not None:
                    existing_spids.add(int(spid))
                break
            parent = parent.getparent()

    # 次のcTn idを算出
    max_id = _get_max_ctn_id(timing_elem)

    for shape_id in audio_shape_ids:
        if shape_id in existing_spids:
            continue
        audio_par = _build_audio_par_xml(shape_id, max_id + 1)
        child_tn_lst.append(audio_par)
        max_id += 3  # 各audioノードはcTn idを3つ使う


def _get_ctn_start_delay(ctn: etree._Element) -> int:
    """cTnの開始遅延(ms)を p:stCondLst/p:cond/@delay から取得する。"""
    st_cond = ctn.find(f"{{{_P_NS}}}stCondLst/{{{_P_NS}}}cond")
    if st_cond is None:
        return 0
    delay_val = st_cond.get("delay", "0")
    if delay_val == "indefinite":
        return 0
    try:
        return int(delay_val)
    except ValueError:
        return 0


def _calc_par_end_ms(par_elem: etree._Element, parent_delay: int) -> int:
    """p:par の終了時刻(ms)を再帰的に計算する。

    ネストした p:cTn の開始遅延を累積して正確な終了時刻を求める。
    音声ノード (dur属性なし) は実質カウントされない。
    """
    ctn = par_elem.find(f"{{{_P_NS}}}cTn")
    if ctn is None:
        return parent_delay

    own_delay = _get_ctn_start_delay(ctn)
    total_delay = parent_delay + own_delay

    # この cTn 自身の dur（p:seq や p:par コンテナは通常 dur="indefinite"）
    dur_val = ctn.get("dur")
    own_end = total_delay
    if dur_val and dur_val != "indefinite":
        try:
            raw_dur = int(dur_val)
            repeat_dur_val = ctn.get("repeatDur")
            repeat_count_val = ctn.get("repeatCount")
            if repeat_dur_val and repeat_dur_val != "indefinite":
                effective_dur = int(repeat_dur_val)
            elif repeat_count_val and repeat_count_val != "indefinite":
                effective_dur = int(raw_dur * float(repeat_count_val))
            else:
                effective_dur = raw_dur
            own_end = total_delay + effective_dur
        except ValueError:
            pass

    # 子を再帰処理
    child_tn_lst = ctn.find(f"{{{_P_NS}}}childTnLst")
    if child_tn_lst is None:
        return own_end

    max_child_end = total_delay
    for child in child_tn_lst:
        if child.tag == f"{{{_P_NS}}}par":
            # p:par はネストしたアニメーションコンテナ
            child_end = _calc_par_end_ms(child, total_delay)
            max_child_end = max(max_child_end, child_end)
        else:
            # p:anim, p:set 等の振る舞いノード: p:cBhvr > p:cTn[@dur] を探す
            for sub_ctn in child.iter(f"{{{_P_NS}}}cTn"):
                sub_dur = sub_ctn.get("dur")
                if sub_dur and sub_dur != "indefinite":
                    try:
                        max_child_end = max(
                            max_child_end,
                            total_delay + _get_ctn_start_delay(sub_ctn) + int(sub_dur),
                        )
                    except ValueError:
                        pass

    return max(own_end, max_child_end)


def _get_max_child_animation_dur_ms(main_seq_ctn: etree._Element) -> int:
    """mainSeq の childTnLst 内の既存アニメーション終了時刻の最大値(ms)を返す。

    各 p:par の終了時刻を再帰計算し、evt="onEnd" による連鎖も解決する。
    音声ノード (dur属性なし) は実質カウントされない。
    """
    child_tn_lst = main_seq_ctn.find(f"{{{_P_NS}}}childTnLst")
    if child_tn_lst is None:
        return 0

    par_list = child_tn_lst.findall(f"{{{_P_NS}}}par")
    if not par_list:
        return 0

    # Step 1: cTn id → par index マップを構築
    ctn_id_to_par_idx: dict[str, int] = {}
    for i, par in enumerate(par_list):
        for ctn in par.iter(f"{{{_P_NS}}}cTn"):
            ctn_id = ctn.get("id")
            if ctn_id:
                ctn_id_to_par_idx[ctn_id] = i

    # Step 2: 各 par の内部所要時間を計算（outer stCondLst の遅延を除く）
    par_internal_durs: list[int] = []
    for par in par_list:
        outer_ctn = par.find(f"{{{_P_NS}}}cTn")
        own_delay = _get_ctn_start_delay(outer_ctn) if outer_ctn is not None else 0
        total_end = _calc_par_end_ms(par, 0)
        par_internal_durs.append(max(0, total_end - own_delay))

    # Step 3: evt="onEnd" 依存を解決して開始時刻を計算
    par_starts: dict[int, int] = {}
    _computing: set[int] = set()  # 循環参照検出用

    def get_start(idx: int) -> int:
        if idx in par_starts:
            return par_starts[idx]
        if idx in _computing:
            # 循環参照検出: 独立開始(0ms)として扱い RecursionError を防ぐ
            par_starts[idx] = 0
            return 0

        _computing.add(idx)
        try:
            par = par_list[idx]
            outer_ctn = par.find(f"{{{_P_NS}}}cTn")
            if outer_ctn is None:
                par_starts[idx] = 0
                return 0

            own_delay = _get_ctn_start_delay(outer_ctn)
            st_cond = outer_ctn.find(f"{{{_P_NS}}}stCondLst/{{{_P_NS}}}cond")

            if st_cond is not None:
                evt = st_cond.get("evt")
                if evt == "onClick":
                    # クリックトリガーは無人再生では発動しないためスキップ (-1 = sentinel)
                    par_starts[idx] = -1
                    return -1
                tn_ref = st_cond.find(f"{{{_P_NS}}}tn")
                if evt in ("onEnd", "onBegin") and tn_ref is not None:
                    ref_id = tn_ref.get("val")
                    ref_par_idx = ctn_id_to_par_idx.get(ref_id)
                    if ref_par_idx is not None and ref_par_idx != idx:
                        ref_start = get_start(ref_par_idx)
                        if ref_start >= 0:
                            if evt == "onEnd":
                                # After Previous: 参照アニメーション終了後に開始
                                base = ref_start + par_internal_durs[ref_par_idx]
                            else:
                                # With Previous (onBegin): 参照アニメーション開始と同時
                                base = ref_start
                            par_starts[idx] = base + own_delay
                            return par_starts[idx]

            par_starts[idx] = own_delay
            return own_delay
        finally:
            _computing.discard(idx)

    max_end = 0
    for i in range(len(par_list)):
        start = get_start(i)
        if start < 0:
            continue  # onClick など無人再生でスキップされるアニメーション
        end = start + par_internal_durs[i]
        max_end = max(max_end, end)
    return max_end


def _get_existing_anim_end_ms(slide) -> int:
    """スライドの既存アニメーション終了時刻(ms)を返す。timing がない場合は 0。"""
    timing_elem = slide.element.find(f"{{{_P_NS}}}timing")
    if timing_elem is None:
        return 0
    main_seq_ctn = timing_elem.find(f".//{{{_P_NS}}}cTn[@nodeType='mainSeq']")
    if main_seq_ctn is None:
        return 0
    return _get_max_child_animation_dur_ms(main_seq_ctn)


def _get_max_ctn_id(timing_elem: etree._Element) -> int:
    """timing要素内の最大cTn idを取得する。"""
    max_id = 0
    for ctn in timing_elem.iter(f"{{{_P_NS}}}cTn"):
        id_val = ctn.get("id")
        if id_val is not None:
            max_id = max(max_id, int(id_val))
    return max_id


def _build_audio_par_xml(
    shape_id: int, start_id: int
) -> etree._Element:
    """音声auto-play用のp:parノードを構築する。

    macOS PowerPoint互換のため p:audio ではなく p:cmd を使用する。
    p:audio > p:cMediaNode はmacOS PowerPointで破損扱いされるが、
    p:cmd type="call" cmd="playFrom(0)" は両プラットフォームで動作する。
    """
    xml = f"""<p:par xmlns:p="{_P_NS}">
  <p:cTn id="{start_id}" fill="hold">
    <p:stCondLst>
      <p:cond delay="0"/>
    </p:stCondLst>
    <p:childTnLst>
      <p:par>
        <p:cTn id="{start_id + 1}" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
          <p:stCondLst>
            <p:cond delay="0"/>
          </p:stCondLst>
          <p:childTnLst>
            <p:cmd type="call" cmd="playFrom(0)">
              <p:cBhvr>
                <p:cTn id="{start_id + 2}" dur="1" fill="hold"/>
                <p:tgtEl>
                  <p:spTgt spid="{shape_id}"/>
                </p:tgtEl>
              </p:cBhvr>
            </p:cmd>
          </p:childTnLst>
        </p:cTn>
      </p:par>
    </p:childTnLst>
  </p:cTn>
</p:par>"""
    return etree.fromstring(xml)


def _build_timing_xml(
    audio_shape_ids: list[int],
    *,
    main_seq_dur_ms: int = 0,
) -> etree._Element:
    """音声auto-play用の完全なp:timing要素を構築する。

    Args:
        audio_shape_ids: 音声シェイプのIDリスト
        main_seq_dur_ms: mainSeq の dur 属性値（ミリ秒）。
            0 の場合は "indefinite"。正の値の場合はその値を使用する。
            LibreOffice互換性のため音声長を渡すと自動送りが正しく動作する。
    """
    audio_pars = ""
    ctn_id = 3
    for shape_id in audio_shape_ids:
        audio_pars += f"""
                <p:par>
                  <p:cTn id="{ctn_id}" fill="hold">
                    <p:stCondLst>
                      <p:cond delay="0"/>
                    </p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="{ctn_id + 1}" presetID="1" presetClass="mediacall" presetSubtype="0" fill="hold" grpId="0" nodeType="afterEffect">
                          <p:stCondLst>
                            <p:cond delay="0"/>
                          </p:stCondLst>
                          <p:childTnLst>
                            <p:cmd type="call" cmd="playFrom(0)">
                              <p:cBhvr>
                                <p:cTn id="{ctn_id + 2}" dur="1" fill="hold"/>
                                <p:tgtEl>
                                  <p:spTgt spid="{shape_id}"/>
                                </p:tgtEl>
                              </p:cBhvr>
                            </p:cmd>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>"""
        ctn_id += 3

    main_seq_dur = str(main_seq_dur_ms) if main_seq_dur_ms > 0 else "indefinite"

    timing_xml = f"""<p:timing xmlns:p="{_P_NS}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="{main_seq_dur}" nodeType="mainSeq">
              <p:childTnLst>{audio_pars}
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0">
                <p:tgtEl><p:sldTgt/></p:tgtEl>
              </p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0">
                <p:tgtEl><p:sldTgt/></p:tgtEl>
              </p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
    return etree.fromstring(timing_xml)


def _estimate_note_duration_ms(slide, *, min_ms: int = _MIN_NOTE_DURATION_MS) -> int:
    """スピーカーノートの文字数から発話時間を推定する（ミリ秒）。

    日本語の平均的な発話速度（約300文字/分）を基準に算出する。
    ノートが空の場合は0を返す。

    Args:
        slide: python-pptxのスライドオブジェクト
        min_ms: 推定結果の最低値（ミリ秒）
    """
    if not slide.has_notes_slide:
        return 0

    try:
        note_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        return 0

    stripped = note_text.strip()
    if not stripped:
        return 0

    duration_ms = int(_estimate_reading_time_ms(stripped))
    return max(min_ms, min(duration_ms, _MAX_NOTE_DURATION_MS))


def _estimate_reading_time_ms(text: str) -> float:
    """テキストの発話時間を推定する（ミリ秒）。

    CJK文字（日本語・中国語・韓国語）とLatin文字で異なるレートを適用:
    - CJK: 約5文字/秒（300文字/分）
    - Latin: 約15文字/秒（単語あたり約5文字 × 180wpm = 900文字/分）
    """
    cjk_count = 0
    latin_count = 0
    for ch in text:
        cp = ord(ch)
        if (
            0x3000 <= cp <= 0x9FFF      # CJK統合漢字 + ひらがな + カタカナ
            or 0xAC00 <= cp <= 0xD7AF   # ハングル音節
            or 0xF900 <= cp <= 0xFAFF   # CJK互換漢字
            or 0xFF00 <= cp <= 0xFFEF   # 全角記号
        ):
            cjk_count += 1
        elif ch.isalnum():
            latin_count += 1
        # 記号・空白は計算に含めない

    cjk_rate = 5.0   # CJK文字/秒
    latin_rate = 15.0  # Latin文字/秒

    cjk_time = cjk_count / cjk_rate if cjk_count else 0
    latin_time = latin_count / latin_rate if latin_count else 0

    classified = cjk_count + latin_count
    if classified > 0:
        # 分類できた文字がある場合はCJK+Latinレートで計算
        return (cjk_time + latin_time) * 1000
    else:
        # 記号のみ等、分類できない場合はフォールバック
        return len(text) / _CHARS_PER_SECOND * 1000


def _find_audio_shape_ids(slide) -> list[int]:
    """スライドから全音声シェイプのIDを取得する。

    a:audioFile要素を持つp:pic要素のcNvPr idをリストで返す。
    """
    ids: list[int] = []
    slide_elem = slide.element
    for audio_file in slide_elem.iter(f"{{{_A_NS}}}audioFile"):
        nv_pr = audio_file.getparent()
        if nv_pr is not None:
            nv_pic_pr = nv_pr.getparent()
            if nv_pic_pr is not None:
                c_nv_pr = nv_pic_pr.find(f"{{{_P_NS}}}cNvPr")
                if c_nv_pr is not None:
                    ids.append(int(c_nv_pr.get("id")))
    return ids
