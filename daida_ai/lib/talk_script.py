"""PPTXスピーカーノートの読み書き + TTSスクリプト入出力"""

from __future__ import annotations

import re
from pathlib import Path
from pptx import Presentation

_DELIMITER_FMT = "--- Slide {idx:03d} ---"
_DELIMITER_RE = re.compile(r"^--- Slide (\d{3,}) ---$")


def read_notes(pptx_path: Path) -> list[str]:
    """PPTXファイルから全スライドのスピーカーノートを取得する。

    Returns:
        スライド順のノートテキストリスト（ノートがなければ空文字）

    Raises:
        FileNotFoundError: ファイルが存在しない
    """
    if not pptx_path.exists():
        raise FileNotFoundError(f"File not found: {pptx_path}")

    prs = Presentation(str(pptx_path))
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide:
            text = slide.notes_slide.notes_text_frame.text
            notes.append(text)
        else:
            notes.append("")
    return notes


def write_notes(
    input_path: Path,
    notes: list[str],
    output_path: Path,
) -> None:
    """PPTXファイルのスピーカーノートを上書きする。

    Args:
        input_path: 入力PPTXファイルパス
        notes: スライド順のノートテキストリスト
        output_path: 出力PPTXファイルパス（input_pathと同一でも可）

    Raises:
        FileNotFoundError: ファイルが存在しない
        ValueError: ノート数とスライド数が一致しない
    """
    if not input_path.exists():
        raise FileNotFoundError(f"File not found: {input_path}")

    prs = Presentation(str(input_path))

    if len(notes) != len(prs.slides):
        raise ValueError(
            f"Slide count ({len(prs.slides)}) does not match notes count ({len(notes)})"
        )

    for slide, note_text in zip(prs.slides, notes):
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = note_text

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def export_tts_script(
    notes: list[str],
    output_path: Path,
    *,
    dict_entries: list[tuple[str, str]] | None = None,
) -> Path:
    """ノートリストをTTSスクリプトファイルにエクスポートする。

    ユーザーが読み上げテキストを確認・修正できるよう、
    区切り線ベースのプレーンテキスト形式で出力する。

    Args:
        notes: スライドごとのスピーカーノート
        output_path: 出力ファイルパス
        dict_entries: 読み辞書エントリ（指定時はノートに自動適用）

    Returns:
        出力ファイルパス
    """
    from daida_ai.lib.pronunciation_dict import apply_dict

    output_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    for i, note in enumerate(notes):
        if dict_entries:
            note = apply_dict(note, dict_entries)
        lines.append(_DELIMITER_FMT.format(idx=i))
        lines.append(note)
        lines.append("")  # 区切り線間の空行
    output_path.write_text("\n".join(lines), encoding="utf-8")
    return output_path


def load_tts_script(script_path: Path) -> list[str]:
    """TTSスクリプトファイルからノートリストを読み込む。

    Args:
        script_path: スクリプトファイルパス

    Returns:
        スライドごとのノートテキストリスト

    Raises:
        FileNotFoundError: ファイルが存在しない
        ValueError: 区切り線が見つからない不正フォーマット
    """
    if not script_path.exists():
        raise FileNotFoundError(f"File not found: {script_path}")

    content = script_path.read_text(encoding="utf-8")
    if not content.strip():
        return []

    lines = content.split("\n")
    notes: list[str] = []
    current_lines: list[str] | None = None
    found_delimiter = False

    for line in lines:
        if _DELIMITER_RE.match(line):
            found_delimiter = True
            if current_lines is not None:
                notes.append(_join_note_lines(current_lines))
            current_lines = []
        elif current_lines is not None:
            current_lines.append(line)

    if current_lines is not None:
        notes.append(_join_note_lines(current_lines))

    if not found_delimiter:
        raise ValueError(
            "No delimiter found in script file. "
            "Expected format: '--- Slide NNN ---'"
        )

    return notes


def _join_note_lines(lines: list[str]) -> str:
    """ノート行を結合し、末尾の空行を除去する。"""
    while lines and not lines[-1]:
        lines.pop()
    return "\n".join(lines)
