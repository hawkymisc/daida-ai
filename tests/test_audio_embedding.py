"""TDD: audio_embed.py — 音声ファイルPPTX埋め込みテスト"""

import zipfile
import pytest
from pathlib import Path
from pptx import Presentation
from pptx.opc.package import RT

from daida_ai.lib.slide_spec import SlideSpec, SlideMetadata, Slide
from daida_ai.lib.slide_builder import build_presentation
from daida_ai.lib.audio_embed import embed_audio_to_pptx
from tests.conftest import DUMMY_MP3_BYTES


@pytest.fixture
def pptx_path(tmp_output_dir: Path) -> Path:
    """テスト用3スライドPPTX"""
    spec = SlideSpec(
        metadata=SlideMetadata(title="T", subtitle="S", event="E"),
        slides=[
            Slide(layout="title_slide", title="スライド1"),
            Slide(layout="title_and_content", title="スライド2", body=["a"]),
            Slide(layout="section_header", title="スライド3"),
        ],
    )
    prs = build_presentation(spec)
    path = tmp_output_dir / "test.pptx"
    prs.save(str(path))
    return path


@pytest.fixture
def audio_dir(tmp_output_dir: Path) -> Path:
    """テスト用音声ファイルディレクトリ（ダミーMP3）"""
    d = tmp_output_dir / "audio"
    d.mkdir()
    for i in [0, 2]:
        (d / f"slide_{i:03d}.mp3").write_bytes(DUMMY_MP3_BYTES)
    return d


class TestEmbedAudio:
    """embed_audio_to_pptx() のテスト"""

    def test_音声ファイルが埋め込まれる(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "output.pptx"
        count = embed_audio_to_pptx(pptx_path, audio_dir, output)
        assert count == 2
        assert output.exists()

    def test_音声なしスライドはスキップされる(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        assert len(prs.slides) == 3

    def test_音声ファイルなしのディレクトリは0件(
        self, pptx_path: Path, tmp_output_dir: Path
    ):
        empty_dir = tmp_output_dir / "empty_audio"
        empty_dir.mkdir()
        output = tmp_output_dir / "output.pptx"
        count = embed_audio_to_pptx(pptx_path, empty_dir, output)
        assert count == 0

    def test_出力ファイルが正常に開ける(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)
        prs = Presentation(str(output))
        assert len(prs.slides) == 3

    def test_メディアパーツがZIPに含まれる(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """ZIPアーカイブ内にppt/media/audio_slideXXX.mp3が存在する"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        with zipfile.ZipFile(str(output), "r") as zf:
            names = zf.namelist()

        assert "ppt/media/audio_slide000.mp3" in names
        assert "ppt/media/audio_slide002.mp3" in names
        assert "ppt/media/audio_slide001.mp3" not in names

    def test_スライドにメディアリレーションシップが設定される(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """埋め込み対象スライドにRT.MEDIAリレーションシップがある"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))

        # スライド0にメディアリレーションシップがある
        slide0_rels = prs.slides[0].part.rels
        media_rels_0 = [r for r in slide0_rels.values() if r.reltype == RT.MEDIA]
        assert len(media_rels_0) > 0

        # スライド1にはメディアリレーションシップがない
        slide1_rels = prs.slides[1].part.rels
        media_rels_1 = [r for r in slide1_rels.values() if r.reltype == RT.MEDIA]
        assert len(media_rels_1) == 0

        # スライド2にメディアリレーションシップがある
        slide2_rels = prs.slides[2].part.rels
        media_rels_2 = [r for r in slide2_rels.values() if r.reltype == RT.MEDIA]
        assert len(media_rels_2) > 0

    def test_スライドXMLに音声シェイプが挿入される(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """埋め込みスライドのXMLにp:pic要素（audioFile付き）がある"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))

        # スライド0のXMLにaudioFile参照がある
        slide0_xml = prs.slides[0].element
        ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
        audio_refs = slide0_xml.findall(".//a:audioFile", ns)
        assert len(audio_refs) > 0


class TestPowerPoint互換性:
    """PowerPointが要求するOOXML準拠の音声埋め込み構造を検証"""

    _ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
        "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
    }

    def test_RT_AUDIOリレーションシップが存在する(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """audioFile r:linkはRT.AUDIO型のリレーションシップを参照すべき"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide0_rels = prs.slides[0].part.rels
        audio_rels = [r for r in slide0_rels.values() if r.reltype == RT.AUDIO]
        assert len(audio_rels) > 0, "RT.AUDIOリレーションシップが必要"

    def test_audioFileのr_linkがRT_AUDIOを参照する(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """a:audioFile r:linkの参照先がRT.AUDIOリレーションシップであること"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        audio_file = slide.element.find(".//a:audioFile", self._ns)
        r_link = audio_file.get(f"{{{self._ns['r']}}}link")

        # r_linkが指すリレーションシップの型を確認
        rel = slide.part.rels[r_link]
        assert rel.reltype == RT.AUDIO, (
            f"audioFile r:linkはRT.AUDIOを参照すべき、実際: {rel.reltype}"
        )

    def test_p14_media拡張要素が存在する(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """PowerPoint 2010+はp14:media拡張要素を要求する"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide_xml = prs.slides[0].element
        p14_media = slide_xml.findall(".//p14:media", self._ns)
        assert len(p14_media) > 0, "p14:media拡張要素が必要"

    def test_p14_mediaのr_embedがRT_MEDIAを参照する(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """p14:media r:embedの参照先がRT.MEDIAリレーションシップであること"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        p14_media = slide.element.find(".//p14:media", self._ns)
        r_embed = p14_media.get(f"{{{self._ns['r']}}}embed")

        rel = slide.part.rels[r_embed]
        assert rel.reltype == RT.MEDIA, (
            f"p14:media r:embedはRT.MEDIAを参照すべき、実際: {rel.reltype}"
        )

    def test_blipにr_embed属性がある(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """a:blipにr:embed属性があり空でないこと"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide_xml = prs.slides[0].element
        # 音声シェイプのblipを取得
        pics = slide_xml.findall(".//p:pic", self._ns)
        audio_pic = None
        for pic in pics:
            if pic.find(".//a:audioFile", self._ns) is not None:
                audio_pic = pic
                break
        assert audio_pic is not None

        blip = audio_pic.find(".//a:blip", self._ns)
        r_embed = blip.get(f"{{{self._ns['r']}}}embed")
        assert r_embed, "a:blipにr:embed属性が必要"

    def test_hlinkClickがppaction_mediaを参照する(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """音声シェイプのa:hlinkClickがRT.HYPERLINKでppaction://mediaを参照すること"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        # 音声シェイプ内のhlinkClickを取得
        pics = slide.element.findall(".//p:pic", self._ns)
        audio_pic = None
        for pic in pics:
            if pic.find(".//a:audioFile", self._ns) is not None:
                audio_pic = pic
                break
        assert audio_pic is not None, "音声シェイプが必要"

        hlink = audio_pic.find(".//a:hlinkClick", self._ns)
        assert hlink is not None, "a:hlinkClick要素が必要"

        r_id = hlink.get(f"{{{self._ns['r']}}}id")
        assert r_id and r_id != "", "hlinkClick r:idは空であってはならない"

        action = hlink.get("action")
        assert action == "ppaction://media", (
            f"actionはppaction://mediaであるべき、実際: {action}"
        )

        # リレーションシップがRT.HYPERLINKであること
        rel = slide.part.rels[r_id]
        assert rel.reltype == RT.HYPERLINK, (
            f"hlinkClickはRT.HYPERLINKを参照すべき、実際: {rel.reltype}"
        )


class Testアイコン配置:
    """音声アイコンがスライド右下に配置されることを検証"""

    _ns = {
        "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
        "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    }

    def test_アイコンがスライド右下に配置される(
        self, pptx_path: Path, audio_dir: Path, tmp_output_dir: Path
    ):
        """16:9スライドでもアイコンが右下隅にあること"""
        output = tmp_output_dir / "output.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, output)

        prs = Presentation(str(output))
        slide = prs.slides[0]
        slide_w = int(prs.slide_width)
        slide_h = int(prs.slide_height)

        pics = slide.element.findall(".//p:pic", self._ns)
        audio_pic = None
        for pic in pics:
            if pic.find(".//a:audioFile", self._ns) is not None:
                audio_pic = pic
                break
        assert audio_pic is not None

        off = audio_pic.find(".//a:off", self._ns)
        assert off is not None, "a:off要素が見つかりません"
        x = int(off.get("x"))
        y = int(off.get("y"))

        # アイコンはスライド右端から1インチ(914400 EMU)以内にあるべき
        assert x > slide_w - 914400, (
            f"x={x} はスライド右端({slide_w})から遠すぎる"
        )
        # アイコンはスライド下端から1インチ以内にあるべき
        assert y > slide_h - 914400, (
            f"y={y} はスライド下端({slide_h})から遠すぎる"
        )
