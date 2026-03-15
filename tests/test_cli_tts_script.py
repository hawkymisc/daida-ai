"""TDD: CLIスクリプトのTTSスクリプトエクスポート/インポート統合テスト

export_tts_script.py と synthesize_audio.py --script の動作を検証する。
PPTXファイルを使い、エクスポート→修正→合成の一連のフローをテストする。
"""

import pytest
from pathlib import Path
from unittest.mock import patch, AsyncMock

from daida_ai.lib.talk_script import (
    read_notes,
    write_notes,
    export_tts_script,
    load_tts_script,
)


@pytest.fixture
def pptx_with_notes(tmp_output_dir: Path) -> Path:
    """スピーカーノート付きのテスト用PPTXファイル"""
    from pptx import Presentation

    prs = Presentation()
    layout = prs.slide_layouts[6]  # blank

    for text in ["生成AIの話", "Claudeの紹介", ""]:
        slide = prs.slides.add_slide(layout)
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = text

    path = tmp_output_dir / "test.pptx"
    prs.save(str(path))
    return path


class Testエクスポートインポート統合:
    """PPTX → エクスポート → 修正 → インポート → 合成 の一連フロー"""

    def test_PPTXからスクリプトをエクスポートして読み戻せる(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"

        export_tts_script(notes, script_path)
        loaded = load_tts_script(script_path)

        assert loaded == notes

    def test_スクリプトを修正して読み戻すと修正が反映される(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"

        export_tts_script(notes, script_path)

        # ユーザーが読みを修正するシミュレーション
        content = script_path.read_text(encoding="utf-8")
        content = content.replace("生成AI", "せいせいエーアイ")
        content = content.replace("Claude", "クロード")
        script_path.write_text(content, encoding="utf-8")

        loaded = load_tts_script(script_path)

        assert loaded[0] == "せいせいエーアイの話"
        assert loaded[1] == "クロードの紹介"
        assert loaded[2] == ""  # 空ノートは空のまま

    @pytest.mark.asyncio
    async def test_スクリプトファイルからの音声合成(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        from daida_ai.lib.synthesize import synthesize_notes
        from tests.conftest import DUMMY_MP3_BYTES_SHORT

        # エクスポート
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"
        export_tts_script(notes, script_path)

        # 修正
        content = script_path.read_text(encoding="utf-8")
        content = content.replace("生成AI", "せいせいエーアイ")
        script_path.write_text(content, encoding="utf-8")

        # スクリプトから読み込んで合成
        modified_notes = load_tts_script(script_path)
        audio_dir = tmp_output_dir / "audio"

        engine = AsyncMock()

        async def fake_synthesize(text, output_path, voice=None):
            output_path.parent.mkdir(parents=True, exist_ok=True)
            output_path.write_bytes(DUMMY_MP3_BYTES_SHORT)
            return output_path

        engine.synthesize.side_effect = fake_synthesize

        results = await synthesize_notes(
            modified_notes, audio_dir, engine=engine
        )

        # 修正済みテキストがエンジンに渡されたことを確認
        first_call_text = engine.synthesize.call_args_list[0][0][0]
        assert "せいせいエーアイ" in first_call_text

        # 非空ノートのみ合成される
        assert results[0] is not None
        assert results[1] is not None
        assert results[2] is None  # 空ノートはスキップ


class Test辞書適用統合:
    """読み辞書を適用したエクスポートの統合テスト"""

    def test_辞書適用でエクスポート時に自動置換される(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"
        dict_entries = [("生成", "せいせい"), ("Claude", "クロード")]

        export_tts_script(notes, script_path, dict_entries=dict_entries)
        loaded = load_tts_script(script_path)

        assert loaded[0] == "せいせいAIの話"
        assert loaded[1] == "クロードの紹介"
        assert loaded[2] == ""  # 空ノートは空のまま

    def test_辞書なしなら従来通りエクスポートされる(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"

        export_tts_script(notes, script_path, dict_entries=None)
        loaded = load_tts_script(script_path)

        assert loaded == notes

    def test_辞書ファイルからロードして適用する一連フロー(
        self, pptx_with_notes: Path, tmp_output_dir: Path
    ):
        from daida_ai.lib.pronunciation_dict import load_dict

        # 辞書ファイル作成
        dict_path = tmp_output_dir / "dict.tsv"
        dict_path.write_text(
            "# 読み辞書\n"
            "生成\tせいせい\n"
            "Claude\tクロード\n",
            encoding="utf-8",
        )

        # ロード → エクスポート → 確認
        entries = load_dict(dict_path)
        notes = read_notes(pptx_with_notes)
        script_path = tmp_output_dir / "tts_script.txt"

        export_tts_script(notes, script_path, dict_entries=entries)
        loaded = load_tts_script(script_path)

        assert "せいせい" in loaded[0]
        assert "クロード" in loaded[1]
