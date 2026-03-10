"""E2Eテスト: テーマ → アウトライン → JSON → PPTX → ノート → 音声埋め込み → スライドショー

全パイプラインをプログラム的に通して動作検証する。
音声合成（edge-tts/VOICEVOX）はネットワーク依存のためモックする。
"""

import json
import pytest
from pathlib import Path
from unittest.mock import AsyncMock, patch
from pptx import Presentation

from daida_ai.lib.outline_parser import parse_outline
from daida_ai.lib.slide_spec import validate_slide_spec, save_slide_spec, load_slide_spec
from daida_ai.lib.slide_builder import build_presentation
from daida_ai.lib.talk_script import read_notes, write_notes
from daida_ai.lib.audio_embed import embed_audio_to_pptx
from daida_ai.lib.slideshow import configure_slideshow


# ── Step 1 相当: Markdownアウトライン ──

SAMPLE_OUTLINE = """\
# AI時代の開発ワークフロー

## 導入: 開発の進化
- 手動コーディングの時代
- IDEとコード補完
- AI支援開発の登場

## 本題1: 生産性の変化
- コード生成速度3倍
- バグ発見率50%向上
- レビュー時間の短縮

## 本題2: 実践的な使い方
- プロンプト設計のコツ
- テスト駆動との組み合わせ
- チーム導入のステップ

## まとめ
- 今日のポイント
- 明日から始める3つのこと
"""

# ── Step 1.5 相当: スライド仕様JSON ──

SAMPLE_SPEC = {
    "metadata": {
        "title": "AI時代の開発ワークフロー",
        "subtitle": "テスト太郎",
        "event": "Tech LT #99",
        "template": "tech",
    },
    "slides": [
        {
            "layout": "title_slide",
            "title": "AI時代の開発ワークフロー",
            "subtitle": "2026/03/10 @ Tech LT #99 - テスト太郎",
        },
        {
            "layout": "section_header",
            "title": "開発の進化",
        },
        {
            "layout": "title_and_content",
            "title": "開発スタイルの変遷",
            "body": [
                "2000年代: 手動コーディング — テキストエディタ中心",
                "2010年代: IDE + コード補完 — IntelliSense/Copilot",
                "2020年代: AI支援開発 — Claude Code/Cursor",
            ],
            "note": "まず開発スタイルの変遷を見ていきましょう。",
        },
        {
            "layout": "two_content",
            "title": "Before / After: 生産性比較",
            "left": {
                "heading": "従来の開発",
                "body": ["コード生成: 1x", "バグ発見: 手動", "レビュー: 2時間/PR"],
            },
            "right": {
                "heading": "AI支援開発",
                "body": ["コード生成: 3x", "バグ発見: 自動50%↑", "レビュー: 30分/PR"],
            },
            "note": "具体的な数字で比較してみましょう。",
        },
        {
            "layout": "title_and_content",
            "title": "実践的な使い方",
            "body": [
                "プロンプト設計: 具体的な指示が鍵",
                "TDD連携: テストファースト + AI実装",
                "チーム導入: 小さく始めて成功体験を積む",
            ],
            "note": "では実践的なテクニックを3つ紹介します。",
        },
        {
            "layout": "title_and_content",
            "title": "まとめ",
            "body": [
                "AI支援開発で生産性3倍",
                "TDDとの組み合わせが最強",
                "明日からClaude Codeを使ってみよう",
            ],
            "note": "最後にまとめです。今日のポイントは3つ。",
        },
    ],
}


class TestE2EPipeline:
    """全パイプラインE2Eテスト"""

    def test_step1_アウトラインパース(self):
        """Step 1: Markdownアウトライン → 構造化データ"""
        outline = parse_outline(SAMPLE_OUTLINE)
        assert outline.title == "AI時代の開発ワークフロー"
        assert len(outline.sections) == 4

    def test_step1_5_スライド仕様バリデーション(self):
        """Step 1.5: スライド仕様JSONバリデーション"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        assert spec.metadata.title == "AI時代の開発ワークフロー"
        assert len(spec.slides) == 6

    def test_step1_5_JSON保存と読み込み(self, tmp_output_dir: Path):
        """Step 1.5: JSON保存 → 読み込みラウンドトリップ"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        json_path = tmp_output_dir / "spec.json"
        save_slide_spec(spec, json_path)

        loaded = load_slide_spec(json_path)
        assert loaded.metadata.title == spec.metadata.title
        assert len(loaded.slides) == len(spec.slides)

    def test_step2_PPTX生成(self, tmp_output_dir: Path):
        """Step 2: スライド仕様JSON → PPTX"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        prs = build_presentation(spec)

        pptx_path = tmp_output_dir / "presentation.pptx"
        prs.save(str(pptx_path))

        assert pptx_path.exists()
        assert pptx_path.stat().st_size > 0

        # 再オープンして検証
        reopened = Presentation(str(pptx_path))
        assert len(reopened.slides) == 6

    def test_step2_スライド内容が正しい(self):
        """Step 2: 生成されたスライドの内容を検証"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        prs = build_presentation(spec)

        # タイトルスライド
        assert prs.slides[0].shapes.title.text == "AI時代の開発ワークフロー"

        # セクションヘッダー
        assert prs.slides[1].shapes.title.text == "開発の進化"

        # コンテンツスライド
        assert prs.slides[2].shapes.title.text == "開発スタイルの変遷"

    def test_step2_ノートが設定されている(self):
        """Step 2: スピーカーノートがスライドに含まれる"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        prs = build_presentation(spec)

        # title_slideとsection_headerにはノートなし
        # title_and_contentスライド（idx=2）にノートあり
        notes_text = prs.slides[2].notes_slide.notes_text_frame.text
        assert "変遷を見ていきましょう" in notes_text

    def test_step3_ノート読み書き(self, tmp_output_dir: Path):
        """Step 3: スピーカーノートの読み出し → 更新"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "notes_test.pptx"
        prs.save(str(pptx_path))

        # 読み出し
        notes = read_notes(pptx_path)
        assert len(notes) == 6

        # 更新
        new_notes = [f"更新されたノート{i}" for i in range(6)]
        write_notes(pptx_path, new_notes, pptx_path)

        # 確認
        result = read_notes(pptx_path)
        assert result == new_notes

    def test_step5_音声埋め込み(self, tmp_output_dir: Path):
        """Step 5: ダミー音声をPPTXに埋め込み"""
        spec = validate_slide_spec(SAMPLE_SPEC)
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "audio_test.pptx"
        prs.save(str(pptx_path))

        # ダミー音声ファイル作成
        audio_dir = tmp_output_dir / "audio"
        audio_dir.mkdir()
        for i in [2, 3, 4, 5]:
            (audio_dir / f"slide_{i:03d}.mp3").write_bytes(
                b"\xff\xfb\x90\x00" + b"\x00" * 100
            )

        # 埋め込み
        output = tmp_output_dir / "final.pptx"
        count = embed_audio_to_pptx(pptx_path, audio_dir, output)
        assert count == 4

        # 再オープン確認
        reopened = Presentation(str(output))
        assert len(reopened.slides) == 6

    def test_フルパイプライン(self, tmp_output_dir: Path):
        """テーマ → アウトライン → JSON → PPTX → ノート更新 → 音声埋め込み → スライドショー"""
        # Step 1: アウトラインパース
        outline = parse_outline(SAMPLE_OUTLINE)
        assert outline.title is not None

        # Step 1.5: スライド仕様JSON
        spec = validate_slide_spec(SAMPLE_SPEC)
        json_path = tmp_output_dir / "spec.json"
        save_slide_spec(spec, json_path)
        loaded_spec = load_slide_spec(json_path)

        # Step 2: PPTX生成
        prs = build_presentation(loaded_spec)
        pptx_path = tmp_output_dir / "presentation.pptx"
        prs.save(str(pptx_path))

        # Step 3: ノート更新
        notes = read_notes(pptx_path)
        updated_notes = [
            n + "（更新済み）" if n else "" for n in notes
        ]
        write_notes(pptx_path, updated_notes, pptx_path)

        # Step 5: 音声埋め込み（ダミー）
        audio_dir = tmp_output_dir / "audio"
        audio_dir.mkdir()
        for i in range(6):
            if updated_notes[i]:
                (audio_dir / f"slide_{i:03d}.mp3").write_bytes(
                    b"\xff\xfb\x90\x00" + b"\x00" * 50
                )

        final_path = tmp_output_dir / "final_presentation.pptx"
        count = embed_audio_to_pptx(pptx_path, audio_dir, final_path)
        assert count > 0

        # Step 6: スライドショー自動再生設定
        _ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
        slideshow_path = tmp_output_dir / "slideshow.pptx"
        configure_slideshow(final_path, slideshow_path)

        # 最終検証
        final_prs = Presentation(str(slideshow_path))
        assert len(final_prs.slides) == 6

        # 全スライドにtransition（自動ページ送り）が設定されている
        for slide in final_prs.slides:
            trans = slide.element.find("p:transition", _ns)
            assert trans is not None, "全スライドにtransitionが必要"
            assert trans.get("advClick") == "0"
            assert trans.get("advTm") is not None

        # 音声付きスライドにtiming（音声auto-play）が設定されている
        audio_slides_with_timing = 0
        for slide in final_prs.slides:
            timing = slide.element.find("p:timing", _ns)
            if timing is not None:
                audio_nodes = timing.findall(".//p:audio", _ns)
                if audio_nodes:
                    audio_slides_with_timing += 1
        assert audio_slides_with_timing > 0, "音声付きスライドにtimingが必要"

        # ノートが更新されていることを確認
        final_notes = read_notes(slideshow_path)
        for i, note in enumerate(final_notes):
            if note:
                assert "更新済み" in note
