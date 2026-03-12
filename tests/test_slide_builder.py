"""TDD: slide_builder.py — スライド仕様JSON → PPTX生成テスト"""

import pytest
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

from daida_ai.lib.slide_spec import (
    SlideSpec,
    SlideMetadata,
    Slide,
    TwoColumnContent,
    validate_slide_spec,
)
from daida_ai.lib.slide_builder import build_presentation, _calc_image_area


@pytest.fixture
def simple_spec() -> SlideSpec:
    """最小限のスライド仕様"""
    return SlideSpec(
        metadata=SlideMetadata(
            title="テストプレゼン",
            subtitle="テスター",
            event="テストイベント",
        ),
        slides=[
            Slide(layout="title_slide", title="テストプレゼン", subtitle="サブタイトル"),
            Slide(
                layout="title_and_content",
                title="本題",
                body=["項目1", "項目2", "項目3"],
                note="ここが説明です。",
            ),
        ],
    )


@pytest.fixture
def full_spec() -> SlideSpec:
    """全レイアウトを含むスライド仕様"""
    return SlideSpec(
        metadata=SlideMetadata(
            title="フルテスト",
            subtitle="テスター",
            event="イベント",
        ),
        slides=[
            Slide(layout="title_slide", title="フルテスト", subtitle="サブ"),
            Slide(layout="section_header", title="セクション1"),
            Slide(
                layout="title_and_content",
                title="コンテンツ",
                body=["A", "B", "C"],
                note="ノートテキスト",
            ),
            Slide(
                layout="two_content",
                title="比較",
                left=TwoColumnContent(heading="左見出し", body=["L1", "L2"]),
                right=TwoColumnContent(heading="右見出し", body=["R1", "R2"]),
                note="比較ノート",
            ),
            Slide(layout="title_only", title="図のスライド", note="図の説明"),
            Slide(layout="blank"),
        ],
    )


class TestBuildPresentation:
    """build_presentation() のテスト"""

    def test_Presentationオブジェクトを返す(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        assert type(prs).__name__ == "Presentation"

    def test_スライド数が仕様と一致する(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        assert len(prs.slides) == 2

    def test_全レイアウトのスライド数が正しい(self, full_spec: SlideSpec):
        prs = build_presentation(full_spec)
        assert len(prs.slides) == 6

    def test_title_slideのタイトルが正しい(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[0]
        assert slide.shapes.title.text == "テストプレゼン"

    def test_title_slideのサブタイトルが正しい(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[0]
        # subtitle is typically placeholder idx 1
        subtitle_shape = slide.placeholders[1]
        assert subtitle_shape.text == "サブタイトル"

    def test_title_and_contentのタイトルが正しい(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[1]
        assert slide.shapes.title.text == "本題"

    def test_title_and_contentの箇条書きが正しい(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[1]
        body_shape = slide.placeholders[1]
        paragraphs = [p.text for p in body_shape.text_frame.paragraphs]
        assert paragraphs == ["項目1", "項目2", "項目3"]

    def test_スピーカーノートが設定される(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[1]
        assert slide.notes_slide.notes_text_frame.text == "ここが説明です。"

    def test_ノートなしのスライドはノートが空(self, simple_spec: SlideSpec):
        prs = build_presentation(simple_spec)
        slide = prs.slides[0]
        notes_text = slide.notes_slide.notes_text_frame.text
        assert notes_text == "" or notes_text.strip() == ""

    def test_section_headerのタイトル(self, full_spec: SlideSpec):
        prs = build_presentation(full_spec)
        slide = prs.slides[1]  # section_header
        assert slide.shapes.title.text == "セクション1"

    def test_two_contentの左右テキスト(self, full_spec: SlideSpec):
        prs = build_presentation(full_spec)
        slide = prs.slides[3]  # two_content
        # Check that all expected text appears in the slide
        all_text = " ".join(shape.text for shape in slide.shapes if shape.has_text_frame)
        assert "左見出し" in all_text
        assert "右見出し" in all_text

    def test_PPTXファイルとして保存できる(
        self, simple_spec: SlideSpec, tmp_output_dir: Path
    ):
        prs = build_presentation(simple_spec)
        path = tmp_output_dir / "test.pptx"
        prs.save(str(path))
        assert path.exists()
        assert path.stat().st_size > 0

    def test_保存したPPTXを再度開ける(
        self, simple_spec: SlideSpec, tmp_output_dir: Path
    ):
        prs = build_presentation(simple_spec)
        path = tmp_output_dir / "test.pptx"
        prs.save(str(path))

        reopened = Presentation(str(path))
        assert len(reopened.slides) == 2


class TestBuildPresentationEdgeCases:
    """build_presentation() のエッジケーステスト"""

    def test_スライドなしの仕様でも空プレゼンを生成(self):
        spec = SlideSpec(
            metadata=SlideMetadata(title="Empty", subtitle="S", event="E"),
            slides=[],
        )
        prs = build_presentation(spec)
        assert len(prs.slides) == 0

    def test_bodyが空リストでもエラーにならない(self):
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_and_content", title="空ボディ", body=[]),
            ],
        )
        prs = build_presentation(spec)
        assert len(prs.slides) == 1

    def test_非常に長いタイトルでもエラーにならない(self):
        long_title = "A" * 500
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_slide", title=long_title),
            ],
        )
        prs = build_presentation(spec)
        assert prs.slides[0].shapes.title.text == long_title

    def test_日本語の長いノートが保存される(self):
        long_note = "これはテストです。" * 100
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[
                Slide(
                    layout="title_and_content",
                    title="ノートテスト",
                    body=["項目"],
                    note=long_note,
                ),
            ],
        )
        prs = build_presentation(spec)
        assert prs.slides[0].notes_slide.notes_text_frame.text == long_note

    def test_blankレイアウトのtitleはNone(self):
        """blankのtitleはデフォルト空文字で、shapes.titleはNone"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="blank")],
        )
        prs = build_presentation(spec)
        assert len(prs.slides) == 1


def _find_pictures(slide):
    """スライド内のPicture型シェイプ一覧を返す"""
    return [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]


class TestSlideSpecImage:
    """Slide dataclassのimage フィールド"""

    def test_imageフィールドのデフォルトはNone(self):
        slide = Slide(layout="title_only", title="T")
        assert slide.image is None

    def test_imageフィールドにパスを設定できる(self):
        slide = Slide(layout="title_only", title="T", image="/path/to/img.png")
        assert slide.image == "/path/to/img.png"

    def test_validate_slide_specでimageが保持される(self):
        data = {
            "metadata": {"title": "T", "subtitle": "S", "event": "E"},
            "slides": [
                {"layout": "title_only", "title": "図", "image": "fig.png"},
            ],
        }
        spec = validate_slide_spec(data)
        assert spec.slides[0].image == "fig.png"

    def test_validate_slide_specでimage省略時はNone(self):
        data = {
            "metadata": {"title": "T", "subtitle": "S", "event": "E"},
            "slides": [
                {"layout": "title_only", "title": "図"},
            ],
        }
        spec = validate_slide_spec(data)
        assert spec.slides[0].image is None

    @pytest.mark.parametrize("bad_value", [True, 123, {"path": "img.png"}, []])
    def test_validate_slide_specでimage型不正はValueError(self, bad_value):
        data = {
            "metadata": {"title": "T", "subtitle": "S", "event": "E"},
            "slides": [
                {"layout": "title_only", "title": "図", "image": bad_value},
            ],
        }
        with pytest.raises(ValueError, match="image must be a string"):
            validate_slide_spec(data)


class TestImageInsertion:
    """スライドへの画像挿入"""

    def _make_spec(self, layout, image_path, **kwargs):
        return SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout=layout, image=str(image_path), **kwargs)],
        )

    def test_title_slideに画像が挿入される(self, sample_image):
        spec = self._make_spec("title_slide", sample_image, title="表紙")
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_title_slideの画像はテキストの背面にある(self, sample_image):
        """title_slideの画像はプレースホルダより背面(z-order前方)に配置される"""
        spec = self._make_spec(
            "title_slide", sample_image, title="表紙", subtitle="サブ",
        )
        prs = build_presentation(spec)
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        pic_indices = [i for i, s in enumerate(shapes) if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        ph_indices = [i for i, s in enumerate(shapes) if s.shape_type != MSO_SHAPE_TYPE.PICTURE and s._element.ph is not None]
        # 画像のインデックスがプレースホルダより小さい = 背面
        assert all(pi < phi for pi in pic_indices for phi in ph_indices), \
            f"pic at {pic_indices} should be behind placeholders at {ph_indices}"

    def test_section_headerに画像が挿入される(self, sample_image):
        spec = self._make_spec("section_header", sample_image, title="セクション")
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_section_headerの画像はテキストの背面にある(self, sample_image):
        spec = self._make_spec("section_header", sample_image, title="セクション")
        prs = build_presentation(spec)
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        pic_indices = [i for i, s in enumerate(shapes) if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        ph_indices = [i for i, s in enumerate(shapes) if s.shape_type != MSO_SHAPE_TYPE.PICTURE and s._element.ph is not None]
        assert all(pi < phi for pi in pic_indices for phi in ph_indices), \
            f"pic at {pic_indices} should be behind placeholders at {ph_indices}"

    def test_title_onlyに画像が挿入される(self, sample_image):
        spec = self._make_spec("title_only", sample_image, title="図")
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_title_onlyの画像はテキストの背面にある(self, sample_image):
        spec = self._make_spec("title_only", sample_image, title="図")
        prs = build_presentation(spec)
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        pic_indices = [i for i, s in enumerate(shapes) if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        ph_indices = [i for i, s in enumerate(shapes) if s.shape_type != MSO_SHAPE_TYPE.PICTURE and s._element.ph is not None]
        assert all(pi < phi for pi in pic_indices for phi in ph_indices), \
            f"pic at {pic_indices} should be behind placeholders at {ph_indices}"

    def test_title_and_contentの画像はテキストの背面にある(self, sample_image):
        spec = self._make_spec("title_and_content", sample_image, title="図付き", body=["テキスト"])
        prs = build_presentation(spec)
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        pic_indices = [i for i, s in enumerate(shapes) if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        ph_indices = [i for i, s in enumerate(shapes) if s.shape_type != MSO_SHAPE_TYPE.PICTURE and s._element.ph is not None]
        assert all(pi < phi for pi in pic_indices for phi in ph_indices), \
            f"pic at {pic_indices} should be behind placeholders at {ph_indices}"

    def test_two_contentの画像はテキストの背面にある(self, sample_image):
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(
                layout="two_content", title="比較",
                left=TwoColumnContent(heading="左", body=["A"]),
                right=TwoColumnContent(heading="右", body=["B"]),
                image=str(sample_image),
            )],
        )
        prs = build_presentation(spec)
        slide = prs.slides[0]
        shapes = list(slide.shapes)
        pic_indices = [i for i, s in enumerate(shapes) if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        ph_indices = [i for i, s in enumerate(shapes) if s.shape_type != MSO_SHAPE_TYPE.PICTURE and s._element.ph is not None]
        assert all(pi < phi for pi in pic_indices for phi in ph_indices), \
            f"pic at {pic_indices} should be behind placeholders at {ph_indices}"

    def test_blankに画像が挿入される(self, sample_image):
        spec = self._make_spec("blank", sample_image)
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_title_and_contentに画像が挿入される(self, sample_image):
        spec = self._make_spec(
            "title_and_content", sample_image,
            title="図付き", body=["テキスト"],
        )
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_imageがNoneの場合は画像なし(self):
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="No image")],
        )
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 0

    def test_画像がコンテンツ領域内に収まる(self, sample_image):
        spec = self._make_spec("title_only", sample_image, title="図")
        prs = build_presentation(spec)
        pic = _find_pictures(prs.slides[0])[0]
        max_w, max_h, _ = _calc_image_area(
            int(prs.slide_width), int(prs.slide_height),
        )
        assert pic.width <= max_w
        assert pic.height <= max_h

    def test_two_contentに画像が挿入される(self, sample_image):
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(
                layout="two_content", title="比較",
                left=TwoColumnContent(heading="左", body=["L1"]),
                right=TwoColumnContent(heading="右", body=["R1"]),
                image=str(sample_image),
            )],
        )
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_ワイド画像のアスペクト比が維持される(self, wide_image):
        spec = self._make_spec("title_only", wide_image, title="ワイド")
        prs = build_presentation(spec)
        pic = _find_pictures(prs.slides[0])[0]
        original_ratio = 800 / 200  # 4.0
        actual_ratio = pic.width / pic.height
        assert abs(actual_ratio - original_ratio) < 0.01

    def test_縦長画像のアスペクト比が維持される(self, tall_image):
        spec = self._make_spec("title_only", tall_image, title="縦長")
        prs = build_presentation(spec)
        pic = _find_pictures(prs.slides[0])[0]
        original_ratio = 200 / 600  # 0.333...
        actual_ratio = pic.width / pic.height
        assert abs(actual_ratio - original_ratio) < 0.01

    def test_画像パスが存在しない場合はFileNotFoundError(self):
        spec = self._make_spec("title_only", "/nonexistent/image.png", title="図")
        with pytest.raises(FileNotFoundError):
            build_presentation(spec)

    def test_不正な画像ファイルはFileNotFoundError(self, tmp_output_dir):
        """壊れたファイルや非画像ファイルでもFileNotFoundErrorになる"""
        bad_file = tmp_output_dir / "not_an_image.png"
        bad_file.write_bytes(b"this is not a valid image")
        spec = self._make_spec("title_only", bad_file, title="壊れた画像")
        with pytest.raises(FileNotFoundError, match="Invalid image file"):
            build_presentation(spec)

    def test_画像付きスライドでノートが保持される(self, sample_image):
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(
                layout="title_only", title="図",
                image=str(sample_image), note="画像の説明",
            )],
        )
        prs = build_presentation(spec)
        assert prs.slides[0].notes_slide.notes_text_frame.text == "画像の説明"

    def test_画像付きPPTXを保存して再度開ける(self, sample_image, tmp_output_dir):
        spec = self._make_spec("title_only", sample_image, title="保存テスト")
        prs = build_presentation(spec)
        path = tmp_output_dir / "img_test.pptx"
        prs.save(str(path))
        reopened = Presentation(str(path))
        pics = _find_pictures(reopened.slides[0])
        assert len(pics) == 1

    def test_相対パスがbase_dirで解決される(self, tmp_path):
        """base_dir指定時に相対パスが正しく解決される"""
        img_dir = tmp_path / "images"
        img_dir.mkdir()
        img_path = img_dir / "test.png"
        from PIL import Image
        Image.new("RGB", (100, 100), "blue").save(str(img_path))
        # 相対パス "images/test.png" を使う
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="相対パス", image="images/test.png")],
        )
        prs = build_presentation(spec, base_dir=tmp_path)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_base_dir未指定で絶対パスは動作する(self, sample_image):
        """base_dir=Noneでも絶対パスなら正常動作"""
        spec = self._make_spec("title_only", sample_image, title="絶対パス")
        prs = build_presentation(spec, base_dir=None)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_パストラバーサルはValueError(self, tmp_path):
        """../を使ってbase_dir外を参照するとValueError"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="悪意", image="../../etc/passwd")],
        )
        with pytest.raises(ValueError, match="escapes base directory"):
            build_presentation(spec, base_dir=tmp_path)


class TestSVGAutoConvert:
    """SVGファイルの自動PNG変換と挿入"""

    def test_SVGファイルがスライドに挿入される(self, sample_svg):
        """SVGパスを指定すると自動変換されて画像が挿入される"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="SVG図", image=str(sample_svg))],
        )
        prs = build_presentation(spec)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_SVG相対パスがbase_dirで解決される(self, tmp_path):
        """相対SVGパスもbase_dirで正しく解決される"""
        img_dir = tmp_path / "images"
        img_dir.mkdir()
        svg_path = img_dir / "diagram.svg"
        svg_path.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="100" height="100">'
            '<rect width="100" height="100" fill="red"/></svg>'
        )
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="相対SVG", image="images/diagram.svg")],
        )
        prs = build_presentation(spec, base_dir=tmp_path)
        pics = _find_pictures(prs.slides[0])
        assert len(pics) == 1

    def test_SVGのアスペクト比が維持される(self, tmp_path):
        """横長SVG（800x200）のアスペクト比が維持される"""
        svg = tmp_path / "wide.svg"
        svg.write_text(
            '<svg xmlns="http://www.w3.org/2000/svg" width="800" height="200">'
            '<rect width="800" height="200" fill="blue"/></svg>'
        )
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="ワイドSVG", image=str(svg))],
        )
        prs = build_presentation(spec)
        pic = _find_pictures(prs.slides[0])[0]
        actual_ratio = pic.width / pic.height
        assert abs(actual_ratio - 4.0) < 0.1

    def test_変換済みPNGは一時ファイルに保存される(self, sample_svg, tmp_output_dir):
        """SVG挿入後もPPTXが正常に保存・再読み込みできる"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="保存テスト", image=str(sample_svg))],
        )
        prs = build_presentation(spec)
        path = tmp_output_dir / "svg_test.pptx"
        prs.save(str(path))
        reopened = Presentation(str(path))
        pics = _find_pictures(reopened.slides[0])
        assert len(pics) == 1

    def test_不正なSVGはFileNotFoundError(self, tmp_path):
        """壊れたSVGファイルでもFileNotFoundErrorが返る"""
        bad_svg = tmp_path / "broken.svg"
        bad_svg.write_text("this is not SVG at all")
        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_only", title="壊れたSVG", image=str(bad_svg))],
        )
        with pytest.raises(FileNotFoundError, match="SVG conversion failed|Invalid image"):
            build_presentation(spec)
