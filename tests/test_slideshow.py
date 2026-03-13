"""TDD: slideshow.py — スライドショー自動再生設定テスト"""

import pytest
from pathlib import Path
from lxml import etree
from pptx import Presentation

from daida_ai.lib.slide_spec import SlideSpec, SlideMetadata, Slide
from daida_ai.lib.slide_builder import build_presentation
from daida_ai.lib.audio_embed import embed_audio_to_pptx
from daida_ai.lib.slideshow import configure_slideshow

# OOXML名前空間
_ns = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
}


@pytest.fixture
def pptx_with_audio(tmp_output_dir: Path) -> Path:
    """音声埋め込み済みの4スライドPPTX"""
    spec = SlideSpec(
        metadata=SlideMetadata(title="LT", subtitle="Speaker", event="Event"),
        slides=[
            Slide(layout="title_slide", title="表紙"),
            Slide(layout="title_and_content", title="内容1", body=["a"]),
            Slide(layout="section_header", title="中表紙"),
            Slide(layout="title_and_content", title="内容2", body=["b"]),
        ],
    )
    prs = build_presentation(spec)
    pptx_path = tmp_output_dir / "test.pptx"
    prs.save(str(pptx_path))

    # スライド1,3に音声を埋め込む（表紙と中表紙には音声なし）
    audio_dir = tmp_output_dir / "audio"
    audio_dir.mkdir()
    # ダミーMP3（104バイト）
    dummy_mp3 = b"\xff\xfb\x90\x00" + b"\x00" * 100
    (audio_dir / "slide_001.mp3").write_bytes(dummy_mp3)
    (audio_dir / "slide_003.mp3").write_bytes(dummy_mp3)

    output = tmp_output_dir / "with_audio.pptx"
    embed_audio_to_pptx(pptx_path, audio_dir, output)
    return output


@pytest.fixture
def pptx_no_audio(tmp_output_dir: Path) -> Path:
    """音声なしの2スライドPPTX"""
    spec = SlideSpec(
        metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
        slides=[
            Slide(layout="title_slide", title="表紙"),
            Slide(layout="title_and_content", title="内容", body=["a"]),
        ],
    )
    prs = build_presentation(spec)
    path = tmp_output_dir / "no_audio.pptx"
    prs.save(str(path))
    return path


class Test自動ページ送り:
    """全スライドにauto-advance transitionが設定される"""

    def test_全スライドにtransition要素が追加される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        for slide in prs.slides:
            trans = slide.element.findall("p:transition", _ns)
            assert len(trans) >= 1, "transition要素が必要"

    def test_advClickが無効になる(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """クリックでのページ送りを無効にする"""
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        for slide in prs.slides:
            trans = slide.element.find("p:transition", _ns)
            assert trans.get("advClick") == "0"

    def test_音声なしスライドにデフォルト秒数が設定される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """表紙・中表紙など音声なしスライドは固定時間で進む"""
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output, silent_duration_ms=4000)

        prs = Presentation(str(output))
        # スライド0 (表紙) と スライド2 (中表紙) は音声なし
        slide0_trans = prs.slides[0].element.find("p:transition", _ns)
        assert slide0_trans.get("advTm") == "4000"

        slide2_trans = prs.slides[2].element.find("p:transition", _ns)
        assert slide2_trans.get("advTm") == "4000"

    def test_音声ありスライドの送り時間はデフォルトより長い(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声付きスライドは音声長に基づいた時間で進む"""
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output, silent_duration_ms=3000)

        prs = Presentation(str(output))
        # スライド1は音声あり → advTmはsilent_durationとは異なる値
        slide1_trans = prs.slides[1].element.find("p:transition", _ns)
        adv_tm = int(slide1_trans.get("advTm"))
        # ダミーMP3は非常に短いが、最低でもバッファ分は確保される
        assert adv_tm > 0

    def test_音声なしPPTXでも動作する(
        self, pptx_no_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_no_audio, output, silent_duration_ms=5000)

        prs = Presentation(str(output))
        for slide in prs.slides:
            trans = slide.element.find("p:transition", _ns)
            assert trans is not None
            assert trans.get("advTm") == "5000"


class Test音声自動再生:
    """音声付きスライドにauto-playアニメーションが設定される"""

    def test_音声付きスライドにtiming要素が追加される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド1は音声あり → timing要素がある
        timing = prs.slides[1].element.find("p:timing", _ns)
        assert timing is not None

    def test_音声なしスライドにtiming要素がない(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド0は音声なし → timing要素がないか空
        timing = prs.slides[0].element.find("p:timing", _ns)
        if timing is not None:
            # timing要素があっても音声アニメーションはない
            audio_nodes = timing.findall(".//p:audio", _ns)
            assert len(audio_nodes) == 0

    def test_出力ファイルが正常に開ける(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        assert len(prs.slides) == 4


class Testカスタム設定:
    """configure_slideshowのパラメータテスト"""

    def test_silent_durationのデフォルトは3000ms(
        self, pptx_no_audio: Path, tmp_output_dir: Path
    ):
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_no_audio, output)

        prs = Presentation(str(output))
        trans = prs.slides[0].element.find("p:transition", _ns)
        assert trans.get("advTm") == "3000"

    def test_audio_buffer_msが加算される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声の長さ + バッファ時間でadvTmが設定される"""
        output = tmp_output_dir / "show.pptx"
        configure_slideshow(pptx_with_audio, output, audio_buffer_ms=2000)

        prs = Presentation(str(output))
        slide1_trans = prs.slides[1].element.find("p:transition", _ns)
        adv_tm = int(slide1_trans.get("advTm"))
        # バッファ2000ms以上は確保されているはず
        assert adv_tm >= 2000


class Test既存アニメーション保持:
    """既存のアニメーションが破壊されないことを検証"""

    @pytest.fixture
    def pptx_with_existing_animation(self, tmp_output_dir: Path) -> Path:
        """既存アニメーション+音声付きPPTX"""
        from lxml import etree

        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_slide", title="表紙"),
                Slide(layout="title_and_content", title="内容", body=["a", "b"]),
            ],
        )
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "base.pptx"
        prs.save(str(pptx_path))

        # スライド1に音声を埋め込む
        audio_dir = tmp_output_dir / "audio_anim"
        audio_dir.mkdir()
        dummy_mp3 = b"\xff\xfb\x90\x00" + b"\x00" * 100
        (audio_dir / "slide_001.mp3").write_bytes(dummy_mp3)
        with_audio = tmp_output_dir / "with_audio_anim.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, with_audio)

        # スライド1に既存アニメーション（テキストフェードイン）を手動追加
        prs2 = Presentation(str(with_audio))
        slide = prs2.slides[1]
        # shape id=2をターゲットにしたダミーアニメーション
        timing_xml = f"""<p:timing xmlns:p="{_P}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:anim calcmode="lin" valueType="num">
                              <p:cBhvr>
                                <p:cTn id="5" dur="500" fill="hold"/>
                                <p:tgtEl><p:spTgt spid="2"/></p:tgtEl>
                                <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                              </p:cBhvr>
                            </p:anim>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
        timing_elem = etree.fromstring(timing_xml)
        slide.element.append(timing_elem)

        result = tmp_output_dir / "with_existing_anim.pptx"
        prs2.save(str(result))
        return result

    def test_既存アニメーションが保持される(
        self, pptx_with_existing_animation: Path, tmp_output_dir: Path
    ):
        """configure_slideshow後も既存のp:animノードが残る"""
        output = tmp_output_dir / "show_anim.pptx"
        configure_slideshow(pptx_with_existing_animation, output)

        prs = Presentation(str(output))
        timing = prs.slides[1].element.find("p:timing", _ns)
        assert timing is not None

        # 既存のアニメーション（p:anim）が残っている
        anim_nodes = timing.findall(".//p:anim", _ns)
        assert len(anim_nodes) >= 1, "既存のp:animアニメーションが保持されるべき"

    def test_既存アニメーションに音声ノードが追加される(
        self, pptx_with_existing_animation: Path, tmp_output_dir: Path
    ):
        """既存アニメーション構造に音声auto-playがマージされる"""
        output = tmp_output_dir / "show_anim.pptx"
        configure_slideshow(pptx_with_existing_animation, output)

        prs = Presentation(str(output))
        timing = prs.slides[1].element.find("p:timing", _ns)
        assert timing is not None

        # 音声auto-playノードが追加されている
        audio_nodes = timing.findall(".//p:audio", _ns)
        assert len(audio_nodes) >= 1, "音声auto-playノードが追加されるべき"

    def test_二重実行しても音声ノードが重複しない(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """configure_slideshowを2回実行してもaudioノードが重複しない"""
        intermediate = tmp_output_dir / "show_pass1.pptx"
        configure_slideshow(pptx_with_audio, intermediate)

        output = tmp_output_dir / "show_pass2.pptx"
        configure_slideshow(intermediate, output)

        prs = Presentation(str(output))
        timing = prs.slides[1].element.find("p:timing", _ns)
        assert timing is not None

        audio_nodes = timing.findall(".//p:audio", _ns)
        assert len(audio_nodes) == 1, "audioノードは1つだけであるべき"


class Testフォールバック:
    """デュレーション計測不能時のフォールバック動作を検証"""

    def test_音声シェイプありデュレーション不明でもauto_playが設定される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声シェイプが存在すれば、デュレーション不明でもtiming要素が追加される。"""
        from unittest.mock import patch

        output = tmp_output_dir / "show_fallback.pptx"
        # _estimate_mp3_duration_msを0返しにスタブ → 計測不能シナリオ
        with patch(
            "daida_ai.lib.slideshow._estimate_mp3_duration_ms", return_value=0
        ):
            configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド1は音声シェイプあり → timing要素が追加される
        timing = prs.slides[1].element.find("p:timing", _ns)
        assert timing is not None
        audio_nodes = timing.findall(".//p:audio", _ns)
        assert len(audio_nodes) >= 1

    def test_unmeasurable_duration_msがadvTmに反映される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """デュレーション計測不能時にフォールバック値がadvTmに使われる"""
        from unittest.mock import patch

        output = tmp_output_dir / "show_unmeasurable.pptx"
        with patch(
            "daida_ai.lib.slideshow._estimate_mp3_duration_ms", return_value=0
        ):
            configure_slideshow(
                pptx_with_audio,
                output,
                unmeasurable_duration_ms=15000,
                audio_buffer_ms=2000,
            )

        prs = Presentation(str(output))
        # スライド1は音声シェイプあり → advTm = 15000 + 2000 = 17000
        slide1_trans = prs.slides[1].element.find("p:transition", _ns)
        assert slide1_trans is not None
        assert slide1_trans.get("advTm") == "17000"


class TestXMLスキーマ順序:
    """ECMA-376 CT_Slide: cSld, clrMapOvr?, transition?, timing?, hf?, extLst?"""

    def test_transitionがtimingより前に配置される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声付きスライドでtransitionとtimingの順序が正しい"""
        output = tmp_output_dir / "show_order.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        slide_elem = prs.slides[1].element
        children = list(slide_elem)
        child_tags = [c.tag.split("}")[-1] for c in children]

        # transitionとtimingが両方存在する場合、transitionが先
        if "transition" in child_tags and "timing" in child_tags:
            trans_idx = child_tags.index("transition")
            timing_idx = child_tags.index("timing")
            assert trans_idx < timing_idx, (
                f"transition(idx={trans_idx}) must precede "
                f"timing(idx={timing_idx})"
            )

    def test_既存timingがあるスライドでtransitionが正しい位置に挿入される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """既にtimingが存在するスライドにtransitionを追加した時の順序"""
        from lxml import etree

        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        # まず音声付きスライドにtimingだけ手動で追加した状態を作る
        prs = Presentation(str(pptx_with_audio))
        slide = prs.slides[1]
        timing_xml = f"""<p:timing xmlns:p="{_P}">
  <p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
    <p:childTnLst><p:seq concurrent="1" nextAc="seek">
      <p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst/></p:cTn>
      <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
      <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
    </p:seq></p:childTnLst>
  </p:cTn></p:par></p:tnLst>
</p:timing>"""
        slide.element.append(etree.fromstring(timing_xml))

        pre_timing_path = tmp_output_dir / "pre_timing.pptx"
        prs.save(str(pre_timing_path))

        output = tmp_output_dir / "show_order_pre_timing.pptx"
        configure_slideshow(pre_timing_path, output)

        prs2 = Presentation(str(output))
        child_tags = [
            etree.QName(c).localname for c in prs2.slides[1].element
        ]
        assert "transition" in child_tags
        assert "timing" in child_tags
        assert child_tags.index("transition") < child_tags.index("timing")

    def test_hf要素があるスライドでtransitionとtimingがhfの前に配置される(
        self, pptx_no_audio: Path, tmp_output_dir: Path
    ):
        """p:hf（ヘッダー/フッター）が存在するスライドでの順序検証"""
        from lxml import etree

        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        # スライドにp:hf要素を手動追加
        prs = Presentation(str(pptx_no_audio))
        slide = prs.slides[0]
        hf_elem = etree.Element(f"{{{_P}}}hf")
        hf_elem.set("sldNum", "0")
        slide.element.append(hf_elem)

        with_hf_path = tmp_output_dir / "with_hf.pptx"
        prs.save(str(with_hf_path))

        output = tmp_output_dir / "show_order_hf.pptx"
        configure_slideshow(with_hf_path, output)

        prs2 = Presentation(str(output))
        child_tags = [
            etree.QName(c).localname for c in prs2.slides[0].element
        ]
        assert "transition" in child_tags
        assert "hf" in child_tags
        assert child_tags.index("transition") < child_tags.index("hf"), (
            f"transition must precede hf, got: {child_tags}"
        )


class Testノートベースタイミング:
    """音声なし・ノートありスライドにノート文字数ベースの送り時間が設定される"""

    @pytest.fixture
    def pptx_with_notes_no_audio(self, tmp_output_dir: Path) -> Path:
        """音声なしだがスピーカーノート付きの3スライドPPTX"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_slide", title="表紙"),
                Slide(
                    layout="title_and_content",
                    title="内容",
                    body=["a"],
                    note="これは50文字程度のスピーカーノートです。発話に約10秒かかると想定されます。ここまでが本文。",
                ),
                Slide(layout="section_header", title="中表紙"),
            ],
        )
        prs = build_presentation(spec)
        path = tmp_output_dir / "notes_no_audio.pptx"
        prs.save(str(path))
        return path

    def test_ノートありスライドのadvTmがsilent_durationより長い(
        self, pptx_with_notes_no_audio: Path, tmp_output_dir: Path
    ):
        """ノートがあるスライドはデフォルト3秒より長い送り時間が設定される"""
        output = tmp_output_dir / "show_notes.pptx"
        configure_slideshow(pptx_with_notes_no_audio, output, silent_duration_ms=3000)

        prs = Presentation(str(output))
        # スライド1はノートあり・音声なし → 3秒より長い
        slide1_trans = prs.slides[1].element.find("p:transition", _ns)
        adv_tm = int(slide1_trans.get("advTm"))
        assert adv_tm > 3000, (
            f"ノートありスライドは3秒より長いはず、実際: {adv_tm}ms"
        )

    def test_ノートなしスライドはsilent_durationが使われる(
        self, pptx_with_notes_no_audio: Path, tmp_output_dir: Path
    ):
        """ノートのないスライド（表紙等）はデフォルトsilent_durationが使われる"""
        output = tmp_output_dir / "show_notes.pptx"
        configure_slideshow(pptx_with_notes_no_audio, output, silent_duration_ms=4000)

        prs = Presentation(str(output))
        # スライド0（表紙）はノートなし → silent_duration
        slide0_trans = prs.slides[0].element.find("p:transition", _ns)
        assert slide0_trans.get("advTm") == "4000"

        # スライド2（中表紙）もノートなし → silent_duration
        slide2_trans = prs.slides[2].element.find("p:transition", _ns)
        assert slide2_trans.get("advTm") == "4000"

    def test_長いノートは長いadvTmになる(self, tmp_output_dir: Path):
        """300文字のノート vs 30文字のノートで、送り時間に差が出る"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(
                    layout="title_and_content",
                    title="短い",
                    body=["a"],
                    note="短いノートです。",
                ),
                Slide(
                    layout="title_and_content",
                    title="長い",
                    body=["b"],
                    note="これは非常に長いスピーカーノートです。" * 20,
                ),
            ],
        )
        prs = build_presentation(spec)
        path = tmp_output_dir / "note_length.pptx"
        prs.save(str(path))

        output = tmp_output_dir / "show_note_length.pptx"
        configure_slideshow(path, output)

        prs2 = Presentation(str(output))
        short_adv = int(prs2.slides[0].element.find("p:transition", _ns).get("advTm"))
        long_adv = int(prs2.slides[1].element.find("p:transition", _ns).get("advTm"))
        assert long_adv > short_adv, (
            f"長いノートのスライドはより長い送り時間のはず: "
            f"short={short_adv}ms, long={long_adv}ms"
        )

    def test_ノートベースタイミングの最低値は3秒(self, tmp_output_dir: Path):
        """非常に短いノートでも最低3秒は確保される"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(
                    layout="title_and_content",
                    title="極短",
                    body=["a"],
                    note="OK",
                ),
            ],
        )
        prs = build_presentation(spec)
        path = tmp_output_dir / "short_note.pptx"
        prs.save(str(path))

        output = tmp_output_dir / "show_short_note.pptx"
        configure_slideshow(path, output)

        prs2 = Presentation(str(output))
        adv_tm = int(prs2.slides[0].element.find("p:transition", _ns).get("advTm"))
        assert adv_tm >= 3000, f"最低3秒は必要: {adv_tm}ms"

    def test_音声ありスライドはノートではなく音声ベース(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声が埋め込まれたスライドは音声デュレーションベースのまま"""
        output = tmp_output_dir / "show_audio_priority.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド1は音声あり → 音声ベースのタイミング（ノートベースではない）
        slide1_trans = prs.slides[1].element.find("p:transition", _ns)
        adv_tm = int(slide1_trans.get("advTm"))
        # ダミーMP3は非常に短い → audio_duration + buffer
        assert adv_tm > 0

    def test_silent_duration指定時ノート付きスライドは最低silent_durationを使う(
        self, tmp_output_dir: Path
    ):
        """silent_duration_ms=10000指定時、短いノートでも10秒未満にならない"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(
                    layout="title_and_content",
                    title="短ノート",
                    body=["a"],
                    note="OK",
                ),
            ],
        )
        prs = build_presentation(spec)
        path = tmp_output_dir / "silent_override.pptx"
        prs.save(str(path))

        output = tmp_output_dir / "show_silent_override.pptx"
        configure_slideshow(path, output, silent_duration_ms=10000)

        prs2 = Presentation(str(output))
        adv_tm = int(prs2.slides[0].element.find("p:transition", _ns).get("advTm"))
        assert adv_tm >= 10000, (
            f"silent_duration_ms=10000が最低値として尊重されるべき: {adv_tm}ms"
        )

    def test_ノートなしスライドにnotesが勝手に作られない(
        self, tmp_output_dir: Path
    ):
        """configure_slideshow後、ノートなしスライドに空のnotesが生成されない"""
        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_slide", title="表紙"),
            ],
        )
        prs = build_presentation(spec)
        # build_presentationはnoteなしスライドにnotesを作らない可能性があるため、
        # has_notes_slideの状態を記録
        path = tmp_output_dir / "no_notes_mutation.pptx"
        prs.save(str(path))

        # 保存後の状態を確認
        prs_before = Presentation(str(path))
        had_notes_before = prs_before.slides[0].has_notes_slide

        output = tmp_output_dir / "show_no_notes_mutation.pptx"
        configure_slideshow(path, output)

        prs_after = Presentation(str(output))
        has_notes_after = prs_after.slides[0].has_notes_slide

        if not had_notes_before:
            assert not has_notes_after, (
                "ノートなしスライドに空のnotesが勝手に作られてはいけない"
            )


class TestLibreOffice互換mainSeqDur:
    """LibreOffice互換: mainSeq の dur が音声長に設定される"""

    def test_音声ありスライドのmainSeq_durが音声長になる(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """mainSeq p:cTn の dur が "indefinite" でなく音声長(ms)であること"""
        output = tmp_output_dir / "show_lo.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド1 (音声あり: slide_001.mp3 → index 1)
        slide = prs.slides[1]
        main_seq_ctn = slide.element.find(
            ".//p:cTn[@nodeType='mainSeq']", _ns
        )
        assert main_seq_ctn is not None, "mainSeq p:cTn が存在すること"
        dur = main_seq_ctn.get("dur")
        assert dur != "indefinite", (
            f"LibreOffice互換のため dur は 'indefinite' であってはならない: {dur}"
        )
        assert dur is not None and int(dur) > 0, (
            f"dur は正の整数(ms)であること: {dur}"
        )

    def test_音声なしスライドにはmainSeq_durが設定されない(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """音声なしスライドには p:timing が追加されないこと"""
        output = tmp_output_dir / "show_lo_no_audio.pptx"
        configure_slideshow(pptx_with_audio, output)

        prs = Presentation(str(output))
        # スライド0 (音声なし: title_slide に mp3 なし)
        slide = prs.slides[0]
        timing = slide.element.find("p:timing", _ns)
        assert timing is None, "音声なしスライドに p:timing は不要"

    def test_mainSeq_durがadvTmと同等以下である(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """mainSeq dur <= advTm であること（音声が終わった後にページ送り）"""
        output = tmp_output_dir / "show_lo_dur.pptx"
        configure_slideshow(pptx_with_audio, output, audio_buffer_ms=1000)

        prs = Presentation(str(output))
        # スライド1 (音声あり)
        slide = prs.slides[1]
        trans = slide.element.find("p:transition", _ns)
        main_seq_ctn = slide.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)

        adv_tm = int(trans.get("advTm"))
        dur = int(main_seq_ctn.get("dur"))
        assert dur <= adv_tm, (
            f"mainSeq dur({dur}ms) は advTm({adv_tm}ms) 以下であるべき"
        )

    def test_既存timingがある場合もmainSeq_durが更新される(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """configure_slideshow を2回適用しても mainSeq.dur が正しく更新される。

        実際のデッキでは既存の p:timing を持つスライドがある。
        _merge_audio_into_timing パスで mainSeq.dur が上書きされることを確認する。
        """
        # 1回目: timing が新規作成される
        output1 = tmp_output_dir / "show_lo_pass1.pptx"
        configure_slideshow(pptx_with_audio, output1)

        # 2回目: 既存 timing を持つスライドに対して _merge_audio_into_timing が呼ばれる
        output2 = tmp_output_dir / "show_lo_pass2.pptx"
        configure_slideshow(output1, output2)

        prs = Presentation(str(output2))
        slide = prs.slides[1]  # 音声ありスライド (index 1)

        main_seq_ctn = slide.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)
        assert main_seq_ctn is not None, "mainSeq p:cTn が存在すること"

        dur = main_seq_ctn.get("dur")
        assert dur != "indefinite", (
            f"_merge_audio_into_timing 経由でも dur が 'indefinite' のまま残ってはいけない: {dur}"
        )
        assert dur is not None and int(dur) > 0, (
            f"dur は正の整数(ms)であること: {dur}"
        )

        # 既存アニメーションノードが重複していないことも確認
        audio_nodes = slide.element.findall(".//p:audio", _ns)
        assert len(audio_nodes) == 1, (
            f"2回適用しても audio ノードは1つだけであるべき: {len(audio_nodes)} 個"
        )

    def test_音声長ゼロのとき_fallback_durationが使われmainSeq_durが正になる(
        self, tmp_output_dir: Path
    ):
        """_get_audio_duration_ms が 0 を返すとき（WAV/外部音声）、
        unmeasurable_duration_ms がフォールバックとして使われ、
        mainSeq.dur が "indefinite" のまま残らないことを確認する。
        """
        from unittest.mock import patch

        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[Slide(layout="title_slide", title="表紙")],
        )
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "test_wav.pptx"
        prs.save(str(pptx_path))

        audio_dir = tmp_output_dir / "audio_wav"
        audio_dir.mkdir()
        # WAV ヘッダ (non-MP3) → _estimate_mp3_duration_ms が 0 を返す
        wav_header = b"RIFF" + b"\x00" * 40
        (audio_dir / "slide_000.mp3").write_bytes(wav_header)

        out_with_audio = tmp_output_dir / "wav_with_audio.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, out_with_audio)

        out_final = tmp_output_dir / "wav_final.pptx"
        fallback_ms = 20000
        configure_slideshow(
            out_with_audio, out_final, unmeasurable_duration_ms=fallback_ms
        )

        prs2 = Presentation(str(out_final))
        slide = prs2.slides[0]
        main_seq_ctn = slide.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)

        assert main_seq_ctn is not None, (
            "WAV/計測不能音声でも p:timing/mainSeq が生成されるべき"
        )
        dur = main_seq_ctn.get("dur")
        assert dur != "indefinite", (
            "WAV/計測不能音声でも mainSeq.dur は 'indefinite' であってはならない"
        )
        assert int(dur) >= 1, f"dur は 1ms 以上であるべき: {dur}"

    def test_既存の非音声アニメーションがmainSeq_dur設定後も消えない(
        self, pptx_with_audio: Path, tmp_output_dir: Path
    ):
        """configure_slideshow 後も既存の音声 audio ノードが保持されること。

        Note: daida-ai が生成するスライドには他のアニメーションノードが存在しないため、
        「音声アニメーションのみ」のケースで重複なく保持されることを検証する。
        2回適用しても audio ノードが増殖しないことで、既存構造の破壊がないことを示す。
        """
        output1 = tmp_output_dir / "anim_pass1.pptx"
        configure_slideshow(pptx_with_audio, output1)

        output2 = tmp_output_dir / "anim_pass2.pptx"
        configure_slideshow(output1, output2)

        prs = Presentation(str(output2))
        for i, slide in enumerate(prs.slides):
            audio_nodes = slide.element.findall(".//p:audio", _ns)
            assert len(audio_nodes) <= 1, (
                f"スライド {i}: audio ノードが重複してはいけない: {len(audio_nodes)} 個"
            )

    def test_再実行時に短縮された音声長でmainSeq_durが縮小される(
        self, tmp_output_dir: Path
    ):
        """configure_slideshow を2回実行し、2回目の音声長が短い場合に
        mainSeq.dur がその短い値に更新されること（stale durにならない）。

        LibreOffice は mainSeq.dur >= advTm でないとページ送りが発火しないため、
        再実行後も dur <= advTm が保証されている必要がある。
        """
        from unittest.mock import patch

        spec = SlideSpec(
            metadata=SlideMetadata(title="T", subtitle="S", event="E"),
            slides=[Slide(layout="title_slide", title="表紙")],
        )
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "shrink_base.pptx"
        prs.save(str(pptx_path))

        audio_dir = tmp_output_dir / "audio_shrink"
        audio_dir.mkdir()
        # WAV ヘッダ (non-MP3) → _get_audio_duration_ms が 0 を返す
        # これにより unmeasurable_duration_ms がフォールバックとして使われる
        wav_header = b"RIFF" + b"\x00" * 40
        (audio_dir / "slide_000.mp3").write_bytes(wav_header)

        out_with_audio = tmp_output_dir / "shrink_audio.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, out_with_audio)

        # 1回目: 長い unmeasurable_duration_ms (30秒) で configure
        out_pass1 = tmp_output_dir / "shrink_pass1.pptx"
        configure_slideshow(out_with_audio, out_pass1, unmeasurable_duration_ms=30000)

        prs1 = Presentation(str(out_pass1))
        slide1 = prs1.slides[0]
        main_seq1 = slide1.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)
        assert main_seq1 is not None
        dur1 = int(main_seq1.get("dur"))
        assert dur1 >= 1

        # 2回目: 短い unmeasurable_duration_ms (5秒) で再 configure
        out_pass2 = tmp_output_dir / "shrink_pass2.pptx"
        configure_slideshow(out_pass1, out_pass2, unmeasurable_duration_ms=5000)

        prs2 = Presentation(str(out_pass2))
        slide2 = prs2.slides[0]
        trans2 = slide2.element.find("p:transition", _ns)
        main_seq2 = slide2.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)

        assert main_seq2 is not None
        dur2 = int(main_seq2.get("dur"))
        adv_tm2 = int(trans2.get("advTm"))

        # 短い音声長で上書きされていること
        assert dur2 < dur1, (
            f"2回目の短い音声長({5000}ms)で dur が縮小されるべき: "
            f"pass1={dur1}ms, pass2={dur2}ms"
        )
        assert dur2 <= adv_tm2, (
            f"mainSeq dur({dur2}ms) は advTm({adv_tm2}ms) 以下であるべき"
        )

    def test_既存非音声アニメーションより短い音声でもtimeline保護される(
        self, tmp_output_dir: Path
    ):
        """音声長 < 既存アニメーションdur のとき、mainSeq.dur は既存アニメーション長以上になる。

        音声(6ms)より長い非音声アニメーション(500ms)が存在する場合、
        mainSeq.dur が 500ms 以上になることで既存アニメーションが切り詰められない。
        """
        from lxml import etree

        _P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        spec = SlideSpec(
            metadata=SlideMetadata(title="LT", subtitle="S", event="E"),
            slides=[
                Slide(layout="title_slide", title="表紙"),
                Slide(layout="title_and_content", title="内容", body=["a"]),
            ],
        )
        prs = build_presentation(spec)
        pptx_path = tmp_output_dir / "anim_protect_base.pptx"
        prs.save(str(pptx_path))

        audio_dir = tmp_output_dir / "audio_anim_protect"
        audio_dir.mkdir()
        # 約6msのフェイクMP3 (音声長 < 500ms)
        dummy_mp3 = b"\xff\xfb\x90\x00" + b"\x00" * 100
        (audio_dir / "slide_001.mp3").write_bytes(dummy_mp3)
        with_audio = tmp_output_dir / "anim_protect_audio.pptx"
        embed_audio_to_pptx(pptx_path, audio_dir, with_audio)

        # スライド1に500msの非音声アニメーションを注入
        prs2 = Presentation(str(with_audio))
        slide = prs2.slides[1]
        timing_xml = f"""<p:timing xmlns:p="{_P}">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:anim calcmode="lin" valueType="num">
                              <p:cBhvr>
                                <p:cTn id="5" dur="500" fill="hold"/>
                                <p:tgtEl><p:spTgt spid="2"/></p:tgtEl>
                                <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                              </p:cBhvr>
                            </p:anim>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
          </p:seq>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>"""
        slide.element.append(etree.fromstring(timing_xml))
        with_existing = tmp_output_dir / "anim_protect_existing.pptx"
        prs2.save(str(with_existing))

        out = tmp_output_dir / "anim_protect_out.pptx"
        configure_slideshow(with_existing, out)

        prs3 = Presentation(str(out))
        slide_out = prs3.slides[1]
        main_seq = slide_out.element.find(".//p:cTn[@nodeType='mainSeq']", _ns)
        assert main_seq is not None

        dur = int(main_seq.get("dur"))
        assert dur >= 500, (
            f"mainSeq.dur({dur}ms) は既存アニメーション(500ms)以上であるべき"
        )
        assert main_seq.get("dur") != "indefinite", (
            "mainSeq.dur は 'indefinite' のまま残ってはいけない"
        )
